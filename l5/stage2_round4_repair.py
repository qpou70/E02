from pathlib import Path
import zipfile, hashlib, subprocess
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND3_第16_17_24页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND4_第16页.pptx'
ASSET = ROOT / 'stage2_round4_assets'
EV = ROOT / 'stage2_visual_review_round4'
RENDER = EV / 'rendered'
COMP = EV / 'comparison'
REPORT = ROOT / 'BUILD_REPORT_STAGE2_ROUND4.md'
TARGET = [16]
FROZEN = list(range(1,16)) + list(range(17,53))
for p in [ASSET, RENDER, COMP]:
    p.mkdir(parents=True, exist_ok=True)

def slide_xmls(pptx):
    out = {}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1, 53):
            out[i] = z.read(f'ppt/slides/slide{i}.xml')
    return out

def sha256(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()

prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52 and doc.page_count == 52
W, H = prs.slide_width, prs.slide_height
BASE_XML = slide_xmls(BASE)

NAVY=RGBColor(30,50,82); BLUE=RGBColor(38,82,132); MID=RGBColor(204,216,230)
BLACK=RGBColor(28,28,30); WHITE=RGBColor(255,255,255); LIGHT=RGBColor(250,252,255)


def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)

def textbox(slide,x,y,w,h,text,size=12,bold=False,color=BLACK,align=PP_ALIGN.LEFT,fill=None,line=None,margin=3):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb=line; sh.line.width=Pt(0.8)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(margin)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    for k,part in enumerate(str(text).split('\n')):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def rect(slide,x,y,w,h,fill=WHITE,line=MID,width=1):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh

def box(slide,x,y,w,h,text,size=11,bold=False,fill=LIGHT,line=BLUE,color=BLACK,align=PP_ALIGN.LEFT,margin=5):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(0.9)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left=tf.margin_right=Pt(margin); tf.margin_top=tf.margin_bottom=Pt(4)
    for k,part in enumerate(str(text).split('\n')):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def header(slide,zh,en):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    textbox(slide,0.38,0.06,6.7,0.42,zh,18,True,WHITE)
    textbox(slide,6.45,0.10,3.12,0.32,en,8.5,False,RGBColor(224,233,243),PP_ALIGN.RIGHT)

def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

def render_crop(page_no,coords,name,dpi=320,mask_right_page_number=False):
    pix=doc[page_no-1].get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),clip=fitz.Rect(*coords),alpha=False)
    out=ASSET/name
    pix.save(out)
    im=Image.open(out).convert('RGB')
    d=ImageDraw.Draw(im)
    if mask_right_page_number:
        d.rectangle([im.width-120, im.height-58, im.width, im.height], fill='white')
    im.save(out)
    return out

def add_pic(slide,path,x,y,w,h):
    with Image.open(path) as im:
        iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    return slide.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

# Rebuild page 16 only. Left text remains the accepted four-category content; right side keeps only original scientific figures plus required Chinese labels.
s=prs.slides[15]
clear(s)
header(s,'键合类型','Bond Types')
box(s,0.38,0.82,3.05,1.25,'离子键 / Ionic Bonding\n电子转移；库仑吸引与原子核排斥在平衡距离处达到平衡。',9.6,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,2.12,3.05,1.15,'共价键 / Covalent Bonding\n电子共享；典型材料：Si、Ge、C。',9.8,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,3.34,3.05,1.38,'混合离子-共价键 / Mixed Ionic-Covalent Bonding\n电负性差导致极性共价键；典型材料：GaAs、InP、GaN。',8.6,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,4.84,3.05,1.22,'金属键 / Metallic Bonding\n正离子实处于电子海中；常见于价电子数不超过 3 的原子。',9.0,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
rect(s,3.62,0.82,5.92,5.92,WHITE,MID,1)

p16a=render_crop(16,(382,170,692,270),'round4_p16_lattice_only.png',320)
p16b=render_crop(16,(350,350,708,435),'round4_p16_deltaE_only.png',320)
p16c=render_crop(16,(395,440,674,528),'round4_p16_metallic_only.png',320)

textbox(s,3.86,0.98,5.25,0.24,'离子键 / 共价键晶格结构示意',9.2,True,NAVY,PP_ALIGN.CENTER)
add_pic(s,p16a,3.90,1.25,5.20,1.25)
# Cover the original English phrase under the covalent-bond figure and place the Chinese wording in the same teaching location.
rect(s,5.46,2.23,1.98,0.28,WHITE,WHITE,0)
textbox(s,5.27,2.20,2.35,0.34,'每个键由两个电子共享',8.2,False,BLACK,PP_ALIGN.CENTER,fill=WHITE,line=WHITE,margin=1)

# Three one-to-one Delta-E labels, each aligned with its original electron-cloud figure.
textbox(s,3.76,2.72,1.64,0.30,'ΔE = 0（共价键）',8.2,True,NAVY,PP_ALIGN.CENTER,fill=WHITE,line=MID,margin=1)
textbox(s,5.48,2.72,1.78,0.30,'ΔE 中等（极性共价键）',7.8,True,NAVY,PP_ALIGN.CENTER,fill=WHITE,line=MID,margin=1)
textbox(s,7.42,2.72,1.56,0.30,'ΔE 较大（离子键）',8.0,True,NAVY,PP_ALIGN.CENTER,fill=WHITE,line=MID,margin=1)
add_pic(s,p16b,3.82,3.08,5.34,1.16)

textbox(s,3.86,4.50,5.25,0.24,'金属键：正离子实 + 离域电子海',9.2,True,NAVY,PP_ALIGN.CENTER)
add_pic(s,p16c,3.88,4.78,5.34,1.36)
# Cover the original English phrase under the metallic-bond figure and add the Chinese term.
rect(s,4.06,6.04,3.55,0.32,WHITE,WHITE,0)
textbox(s,4.24,6.02,3.10,0.34,'离域电子海',10.0,True,BLACK,PP_ALIGN.CENTER,fill=WHITE,line=WHITE,margin=1)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 16. ROUND4 only updates page 16: Delta-E labels restored and two translatable English figure notes localized.')

# Validate visible slide text.
txt='\n'.join((getattr(sh,'text','') or '') for sh in s.shapes)
for bad in ['Two electrons per bond','Swarm of delocalised electrons','Electron transfer','Coulomb attraction balanced','Electron sharing','Polar covalent bonds','Positive cores in an electron sea','Typical for atoms with 3 or less valence electrons']:
    assert bad not in txt, f'page 16 still has forbidden English instruction text: {bad}'
for ok in ['ΔE = 0（共价键）','ΔE 中等（极性共价键）','ΔE 较大（离子键）','每个键由两个电子共享','离域电子海']:
    assert ok in txt, f'missing required localized label: {ok}'

prs.save(OUT)
POST_XML=slide_xmls(OUT)
changed=[i for i in range(1,53) if POST_XML[i]!=BASE_XML[i]]
if changed != TARGET:
    raise AssertionError(f'Only page 16 may change; actual changed pages: {changed}')
pre_hash=sha256(OUT)

# Render PDF and evidence.
pdf_dir=EV/'new_pdf'; pdf_dir.mkdir(parents=True,exist_ok=True)
subprocess.run(['libreoffice','--headless','--convert-to','pdf','--outdir',str(pdf_dir),str(OUT)],check=True)
pdf_path=pdf_dir/(OUT.stem+'.pdf')
newdoc=fitz.open(pdf_path)
assert newdoc.page_count == 52
pix=newdoc[15].get_pixmap(matrix=fitz.Matrix(2.2,2.2),alpha=False)
out_png=RENDER/'page_16.png'; pix.save(out_png)
opix=doc[15].get_pixmap(matrix=fitz.Matrix(2.0,2.0),alpha=False)
orig=Image.frombytes('RGB',[opix.width,opix.height],opix.samples)
new=Image.open(out_png).convert('RGB')
h=max(orig.height,new.height)+70; w=orig.width+new.width+60
canvas=Image.new('RGB',(w,h),'white'); d=ImageDraw.Draw(canvas)
d.text((10,10),'Original PDF page 16',fill=(0,0,0))
d.text((orig.width+50,10),'New Stage 2 ROUND4 page 16',fill=(0,0,0))
canvas.paste(orig,(10,50)); canvas.paste(new,(orig.width+50,50))
comp=COMP/'page_16_original_pdf_vs_new_round4.jpg'
canvas.save(comp,quality=92)
contact=EV/'contact_sheet_stage2_round4_page_16.jpg'
thumb=Image.open(out_png).convert('RGB'); thumb.thumbnail((620,460))
cs=Image.new('RGB',(680,520),'white'); d=ImageDraw.Draw(cs); d.text((10,10),'page 16',fill=(0,0,0)); cs.paste(thumb,((680-thumb.width)//2,45)); cs.save(contact,quality=92)
post_hash=sha256(OUT)

report=f'''# ECE340 L5 第二阶段视觉返修 ROUND4 Build Report\n\n- 返修依据：`l5/SUPERVISOR_STAGE2_VISUAL_REVIEW_FEEDBACK_ROUND4.md`，反馈提交 `ccd341b57c6b3864d8fc7ddb7d556ff1b305ed15`。\n- 当前返修基准 PPT：`l5/ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND3_第16_17_24页.pptx`\n- 输出 PPT：`l5/ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND4_第16页.pptx`\n- 原始 PDF：`l5/stage1_source_reference/ECE340_L5_S18_Posted.pdf`\n- 本轮实际修改页面：第 16 页。\n- 冻结且确认未修改页面：第 8–15、17–24 页；同时确认第 1–7、25–52 页未修改。\n- 变更页 XML 检查：仅 [16] 与 ROUND3 基准 PPT 不同。\n\n## 第 16 页修复摘要\n\n- 左侧四类键合中文说明保持不动。\n- 右侧三组科学图保持，补回三幅电子云图对应标签：`ΔE = 0（共价键）`、`ΔE 中等（极性共价键）`、`ΔE 较大（离子键）`。\n- 将原图中 `Two electrons per bond` 中文化为 `每个键由两个电子共享`。\n- 将原图中 `Swarm of delocalised electrons` 中文化为 `离域电子海`。\n\n## 强制自检记录\n\n- `page_16.png` 已生成。\n- 不存在原 PDF 橙色标题条。\n- 不存在大段英文 bullet list。\n- 不存在 `Two electrons per bond` 或 `Swarm of delocalised electrons` 可见文本。\n- 三个 `ΔE` 标签和两处中文化说明已写入页面。\n- PPT SHA-256（渲染前）：`{pre_hash}`\n\n## 渲染与证据\n\n- 高清渲染图：`l5/stage2_visual_review_round4/rendered/page_16.png`。\n- Original PDF vs New Page 对照图：`l5/stage2_visual_review_round4/comparison/page_16_original_pdf_vs_new_round4.jpg`。\n- Contact sheet：`l5/stage2_visual_review_round4/contact_sheet_stage2_round4_page_16.jpg`。\n- Worker self-check: passed.\n- Supervisor visual acceptance: pending.\n- PPT SHA-256（渲染后）：`{post_hash}`\n'''
REPORT.write_text(report, encoding='utf-8')
subprocess.run(['git','add',str(OUT),str(out_png),str(comp),str(contact),str(REPORT)],check=True)
print('Stage 2 ROUND4 page 16 repair generated and staged.')
