from pathlib import Path
import shutil
import zipfile
import hashlib
import subprocess
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND5_第16页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_最终候选版_第8-24页.pptx'
EV = ROOT / 'final_qa_08_24'
RENDER = EV / 'rendered'
COMP = EV / 'comparison'
PDF_DIR = EV / 'final_pdf'
REPORT = ROOT / 'BUILD_REPORT_FINAL_QA_08_24.md'
PAGES = list(range(8, 25))
for p in [EV, RENDER, COMP, PDF_DIR]:
    p.mkdir(parents=True, exist_ok=True)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def slide_xmls(pptx: Path):
    data = {}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1, 53):
            data[i] = z.read(f'ppt/slides/slide{i}.xml')
    return data


def extract_notes(prs: Presentation, idx: int) -> str:
    try:
        return prs.slides[idx - 1].notes_slide.notes_text_frame.text or ''
    except Exception:
        return ''


def draw_label(draw, xy, text):
    draw.text(xy, text, fill=(0, 0, 0))

assert BASE.exists(), f'Missing ROUND5 base PPT: {BASE}'
assert SRC.exists(), f'Missing source PDF: {SRC}'

# Final Candidate packaging: byte-for-byte copy of ROUND5 base. No slide content edits.
shutil.copyfile(BASE, OUT)

# Regression: slide XML must remain identical for every slide, especially frozen pages 8-24.
base_xml = slide_xmls(BASE)
out_xml = slide_xmls(OUT)
xml_changed = [i for i in range(1, 53) if base_xml[i] != out_xml[i]]
if xml_changed:
    raise AssertionError(f'Final candidate must be a direct package copy; slide XML changed: {xml_changed}')

base_prs = Presentation(BASE)
prs = Presentation(OUT)
assert len(prs.slides) == 52, f'Expected 52 slides, found {len(prs.slides)}'
assert len(base_prs.slides) == 52, f'Expected 52 base slides, found {len(base_prs.slides)}'

# Notes check is a regression check: final candidate must preserve ROUND5 notes exactly.
notes_changed = []
notes_empty = []
notes_missing_sources = []
notes_meta = []
for i in PAGES:
    base_notes = extract_notes(base_prs, i)
    out_notes = extract_notes(prs, i)
    if base_notes != out_notes:
        notes_changed.append(i)
    if not out_notes.strip():
        notes_empty.append(i)
    if '[Sources]' not in out_notes:
        notes_missing_sources.append(i)
    for bad in ['教师应当', '本页建议', '讲授顺序', '制作说明', '施工说明', 'Placeholder', '待替换']:
        if bad in out_notes:
            notes_meta.append((i, bad))

if notes_changed:
    raise AssertionError(f'Notes changed during final packaging; changed pages: {notes_changed}')
if notes_empty or notes_meta:
    raise AssertionError(f'Notes regression found. empty={notes_empty}, meta={notes_meta}')

notes_status = 'passed: final candidate preserves ROUND5 notes byte-for-byte; no notes were cleared or altered during packaging'
notes_sources_detail = f'[Sources] missing in final candidate pages (same as ROUND5 baseline, not introduced by packaging): {notes_missing_sources}'

# Export Final Candidate PDF and render pages 8-24.
subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', str(PDF_DIR), str(OUT)], check=True)
pdf_path = PDF_DIR / (OUT.stem + '.pdf')
assert pdf_path.exists(), f'Missing exported PDF: {pdf_path}'
newdoc = fitz.open(pdf_path)
origdoc = fitz.open(SRC)
assert newdoc.page_count == 52, f'Final PDF has {newdoc.page_count} pages, expected 52'
assert origdoc.page_count == 52, f'Original PDF has {origdoc.page_count} pages, expected 52'

render_paths = []
comparison_paths = []
for page in PAGES:
    pix = newdoc[page - 1].get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
    png_path = RENDER / f'page_{page:02d}.png'
    pix.save(png_path)
    render_paths.append(png_path)

    opix = origdoc[page - 1].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
    npix = newdoc[page - 1].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
    orig = Image.frombytes('RGB', [opix.width, opix.height], opix.samples)
    new = Image.frombytes('RGB', [npix.width, npix.height], npix.samples)
    canvas = Image.new('RGB', (orig.width + new.width + 70, max(orig.height, new.height) + 80), 'white')
    d = ImageDraw.Draw(canvas)
    draw_label(d, (10, 10), f'Original PDF page {page}')
    draw_label(d, (orig.width + 60, 10), f'Final Candidate page {page}')
    canvas.paste(orig, (10, 55))
    canvas.paste(new, (orig.width + 60, 55))
    comp_path = COMP / f'page_{page:02d}_original_pdf_vs_final_candidate.jpg'
    canvas.save(comp_path, quality=92)
    comparison_paths.append(comp_path)

# Overall contact sheet for pages 8-24.
thumbs = []
for page, path in zip(PAGES, render_paths):
    im = Image.open(path).convert('RGB')
    im.thumbnail((360, 235))
    tile = Image.new('RGB', (390, 280), 'white')
    d = ImageDraw.Draw(tile)
    draw_label(d, (10, 8), f'page {page:02d}')
    tile.paste(im, ((390 - im.width) // 2, 38))
    thumbs.append(tile)
cols = 4
rows = 5
sheet = Image.new('RGB', (cols * 390, rows * 280), 'white')
for idx, tile in enumerate(thumbs):
    x = (idx % cols) * 390
    y = (idx // cols) * 280
    sheet.paste(tile, (x, y))
contact_path = EV / 'contact_sheet_final_qa_pages_08_24.jpg'
sheet.save(contact_path, quality=92)

visual_auto_status = 'generated; actual visual inspection by model after artifact download required before final chat report'
mppt_status = 'not performed / pending'

out_hash = sha256(OUT)
report = f'''# ECE340 L5 Final Candidate QA Build Report（第 8–24 页）

- Supervisor instruction source: `l5/SUPERVISOR_STAGE2_ACCEPTED_AND_FINAL_QA_INSTRUCTION.md`, commit `f8c5b32b47afafcdf1d4034ecd0000103f8590de`.
- ROUND5 基准 PPT 路径：`{BASE}`
- 最终候选版路径：`{OUT}`
- 最终候选版 SHA-256：`{out_hash}`
- 52 页页数确认：passed（PPT: {len(prs.slides)} slides; exported PDF: {newdoc.page_count} pages）
- 第 8–24 页冻结确认：passed（最终候选版为 ROUND5 基准复制；slide XML 1–52 全部无变化）
- 本轮实际修改页面内容：无

## 17 张 PNG 路径

''' + '\n'.join(f'- `{p}`' for p in render_paths) + f'''

## 17 张 Original PDF vs Final Candidate 对照图路径

''' + '\n'.join(f'- `{p}`' for p in comparison_paths) + f'''

## Contact sheet

- `{contact_path}`

## Notes / [Sources] 检查结果

- Notes / Sources 检查：{notes_status}
- {notes_sources_detail}
- empty notes pages: {notes_empty}
- notes changed during packaging: {notes_changed}
- meta/施工说明 forbidden-term hits: {notes_meta}

## Visual regression

- Automated generation status: {visual_auto_status}
- Actual 17-page visual inspection by model/user after artifact download: pending at build time

## Microsoft PowerPoint actual open check

- Microsoft PowerPoint actual open check: {mppt_status}

## Supervisor status

- Supervisor visual acceptance: pending
- Supervisor final acceptance: pending
'''
REPORT.write_text(report, encoding='utf-8')

subprocess.run(['git', 'add', str(OUT), str(REPORT), str(EV)], check=True)
print('Final Candidate QA packaging complete')
print('PPT:', OUT)
print('Rendered:', RENDER)
print('Comparison:', COMP)
print('Contact:', contact_path)
print('Report:', REPORT)
