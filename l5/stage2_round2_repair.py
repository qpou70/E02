from pathlib import Path
import zipfile, hashlib
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND1_第8_10_11_14_16_17_19_22_24页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND2_第11_14_16_17_19_22_24页.pptx'
ASSET = ROOT / 'stage2_round2_assets'
REPORT = ROOT / 'BUILD_REPORT_STAGE2_ROUND2.md'
TARGET = [11,14,16,17,19,22,24]
FROZEN = [8,9,10,12,13,15,18,20,21,23]
FORBIDDEN = ['源页','原页裁取','去除标题条','已清理','页面只保留','只保留一处','等比裁取','不叠加公式','真实设备照片','已完成中文化','本页已重建','保留原图','施工','待替换','占位','Placeholder','此处放图']
ASSET.mkdir(exist_ok=True)

prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52 and doc.page_count == 52
W, H = prs.slide_width, prs.slide_height

NAVY=RGBColor(30,50,82); BLUE=RGBColor(38,82,132); TEAL=RGBColor(45,119,122)
GRAY=RGBColor(92,98,108); BLACK=RGBColor(28,28,30); WHITE=RGBColor(255,255,255)
PALE=RGBColor(240,246,252); MID=RGBColor(204,216,230); LIGHT=RGBColor(250,252,255)


def slide_xmls(pptx):
    out={}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1,53):
            out[i]=z.read(f'ppt/slides/slide{i}.xml')
    return out

BASE_XML = slide_xmls(BASE)


def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)


def textbox(slide,x,y,w,h,text,size=12,bold=False,color=BLACK,align=PP_ALIGN.LEFT,fill=None,line=None,margin=3):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb=line; sh.line.width=Pt(0.8)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(margin)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    for k, part in enumerate(str(text).split('\n')):
        if k:
            p=tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh


def rect(slide,x,y,w,h,fill=WHITE,line=MID,width=1):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh


def box(slide,x,y,w,h,text,size=11,bold=False,fill=LIGHT,line=BLUE,color=BLACK,align=PP_ALIGN.LEFT,margin=5):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(0.9)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Pt(margin); tf.margin_top=tf.margin_bottom=Pt(4); tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    for k, part in enumerate(str(text).split('\n')):
        if k:
            p=tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh


def header(slide,zh,en):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    textbox(slide,0.38,0.06,6.45,0.42,zh,18,True,WHITE)
    textbox(slide,6.48,0.10,3.10,0.32,en,8.5,False,RGBColor(224,233,243),PP_ALIGN.RIGHT)


def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text


def render_crop(page_no,coords,name,dpi=300,mask_right_page_number=True,mask_left_px=0,mask_top_px=0,mask_bottom_px=0):
    pix=doc[page_no-1].get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),clip=fitz.Rect(*coords),alpha=False)
    out=ASSET/name; pix.save(out)
    im=Image.open(out).convert('RGB'); d=ImageDraw.Draw(im)
    if mask_right_page_number:
        d.rectangle([im.width-96, im.height-42, im.width, im.height], fill='white')
    if mask_left_px:
        d.rectangle([0,0,mask_left_px,im.height], fill='white')
    if mask_top_px:
        d.rectangle([0,0,im.width,mask_top_px], fill='white')
    if mask_bottom_px:
        d.rectangle([0,im.height-mask_bottom_px,im.width,im.height], fill='white')
    im.save(out)
    return out


def add_pic(slide,path,x,y,w,h):
    with Image.open(path) as im:
        iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    return slide.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))


def simple_table(slide,x,y,cols,rows,widths,row_h=0.42,font=8.6):
    cx=x
    for j,c in enumerate(cols):
        box(slide,cx,y,widths[j],row_h,c,font,True,PALE,BLUE,NAVY,PP_ALIGN.CENTER,margin=2); cx+=widths[j]
    for i,row in enumerate(rows):
        cx=x
        for j,c in enumerate(row):
            fill=RGBColor(252,253,255) if i%2==0 else WHITE
            box(slide,cx,y+row_h*(i+1),widths[j],row_h,str(c),font,False,fill,MID,BLACK,PP_ALIGN.CENTER if j else PP_ALIGN.LEFT,margin=2); cx+=widths[j]

# Page 11: keep one clear reaction equation and clean the bottom repeated half formula
s=prs.slides[10]; clear(s); header(s,'MOCVD：现代系统','MOCVD: Today')
box(s,1.05,0.82,7.90,0.48,'(CH₃)₃Ga + AsH₃  →  GaAs + 3CH₄',15,False,WHITE,BLUE,BLACK,PP_ALIGN.CENTER)
p11=render_crop(11,(36,165,756,472),'round2_p11_mocvd_today_photos.png',300,True,mask_bottom_px=8)
rect(s,0.52,1.46,8.96,5.10,WHITE,MID,1); add_pic(s,p11,0.66,1.62,8.68,4.82)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 11.')

# Page 14: keep only the original visible source already present in the figure crop
s=prs.slides[13]; clear(s); header(s,'生产型 MBE 反应器','Production MBE Reactor')
p14=render_crop(14,(36,126,756,505),'round2_p14_production_mbe.png',300,True)
rect(s,0.52,0.86,8.96,5.90,WHITE,MID,1); add_pic(s,p14,0.66,1.02,8.68,5.50)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 14. Additional equipment source retained only in notes: www.mbe-komponenten.de')

# Page 16: preserve the corrected four categories and restore original teaching diagrams from the PDF body
s=prs.slides[15]; clear(s); header(s,'键合类型','Bond Types')
# Chinese classification column
box(s,0.38,0.82,3.05,1.25,'离子键 / Ionic Bonding\n电子转移；库仑吸引与原子核排斥在平衡距离处达到平衡。',9.6,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,2.12,3.05,1.15,'共价键 / Covalent Bonding\n电子共享；典型材料：Si、Ge、C。',9.8,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,3.34,3.05,1.38,'混合离子-共价键 / Mixed Ionic-Covalent Bonding\n电负性差导致极性共价键；典型材料：GaAs、InP、GaN。',8.6,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,4.84,3.05,1.22,'金属键 / Metallic Bonding\n正离子实处于电子海中；常见于价电子数不超过 3 的原子。',9.0,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
# Original teaching diagrams: full body crop preserves ionic/covalent lattice, electronegativity-polarization sequence, and electron sea diagram.
p16=render_crop(16,(36,118,756,532),'round2_p16_original_bonding_diagrams.png',300,True)
rect(s,3.62,0.82,5.92,5.92,WHITE,MID,1); add_pic(s,p16,3.74,0.96,5.68,5.62)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 16.')

# Page 17: title semantics only, no table redesign
s=prs.slides[16]
found_title = False
for sh in s.shapes:
    txt=getattr(sh,'text','') or ''
    if '元素的电子组态与价电子占据' in txt or '价电子占据' in txt:
        sh.text_frame.clear()
        p=sh.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
        r=p.add_run(); r.text='原子轨道的电子占据'; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(18); r.font.bold=True; r.font.color.rgb=WHITE
        found_title = True
if not found_title:
    # fallback: leave content unchanged but place the faithful title on top band if the previous title was an image/text-free header
    pass
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 17.')

# Page 19: remove residual orange title strip without losing the bottom orbitals
s=prs.slides[18]; clear(s); header(s,'原子轨道波函数的空间形状','Orbital Wave Function Shape')
p19=render_crop(19,(36,132,756,532),'round2_p19_orbitron_full_clean.png',300,True,mask_top_px=10)
rect(s,0.48,0.82,9.04,6.16,WHITE,MID,1); add_pic(s,p19,0.60,0.96,8.80,5.84)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 19; Orbitron orbital image source retained in notes only.')

# Page 22: keep Chinese logic boxes, rebuild the right science figures with a clean figure-only crop
s=prs.slides[21]; clear(s); header(s,'周期势与 E-k 关系','Periodic Potential and E-k Relation')
items=[
    ('定性理解','Streetman 教材与本课程提供的是能带形成的定性图像。'),
    ('理论范围','能带形成的详细理论超出本课程范围。'),
    ('周期势','在晶格周期势中求解薛定谔方程会产生能带。'),
    ('E-k 关系','解给出电子能量 E 与晶体中动量波矢 k 的关系。'),
    ('光学类比','可类比干涉滤光片和蝴蝶翅膀衍射结构，只是作用尺度对应电子波长。')]
y=0.88
for title,body in items:
    box(s,0.45,y,4.55,0.90,f'{title}\n{body}',9.5,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT,margin=5)
    y += 1.05
# Tighter right crop with left margin removed to eliminate broken English body fragments.
p22=render_crop(22,(400,128,756,518),'round2_p22_clean_figures.png',320,True,mask_left_px=28)
rect(s,5.25,0.88,4.22,5.96,WHITE,MID,1); add_pic(s,p22,5.36,1.02,4.00,5.64)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 22.')

# Page 24: one-character semantic correction only
s=prs.slides[23]
for sh in s.shapes:
    txt=getattr(sh,'text','') or ''
    if 'N ↑ Si 原子形成晶体时' in txt or 'N↑ Si 原子形成晶体时' in txt:
        sh.text = txt.replace('N ↑ Si 原子形成晶体时','N 个 Si 原子形成晶体时').replace('N↑ Si 原子形成晶体时','N 个 Si 原子形成晶体时')
    elif 'N ↑' in txt:
        sh.text = txt.replace('N ↑','N 个')
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 24.')

# Validate no maker words on visible slide text and no out-of-bounds content.
for p in TARGET:
    slide=prs.slides[p-1]
    for sh in list(slide.shapes):
        txt = getattr(sh,'text','') or ''
        if any(bad in txt for bad in FORBIDDEN):
            raise AssertionError(f'forbidden maker word on page {p}: {txt}')
        assert sh.left >= -1000 and sh.top >= -1000 and sh.left + sh.width <= W + 1000 and sh.top + sh.height <= H + 1000, f'out of bounds page {p}'

# Explicit semantic checks on slide text.
text17='\n'.join(getattr(sh,'text','') or '' for sh in prs.slides[16].shapes)
assert '原子轨道的电子占据' in text17, 'page 17 title not updated'
text24='\n'.join(getattr(sh,'text','') or '' for sh in prs.slides[23].shapes)
assert 'N ↑' not in text24 and 'N 个 Si 原子形成晶体时' in text24, 'page 24 N title not corrected'

prs.save(OUT)
NEW_XML = slide_xmls(OUT)
changed=[i for i in range(1,53) if NEW_XML[i] != BASE_XML[i]]
assert set(changed).issubset(set(TARGET)), f'unexpected changed pages: {changed}'
for p in FROZEN + list(range(1,8)) + list(range(25,53)):
    assert NEW_XML[p] == BASE_XML[p], f'frozen/outside page changed: {p}'

h=hashlib.sha256(OUT.read_bytes()).hexdigest()
REPORT.write_text(f'''# ECE340 L5 第二阶段视觉返修 ROUND2 Build Report

- 返修依据：`l5/SUPERVISOR_STAGE2_VISUAL_REVIEW_FEEDBACK_ROUND2.md`，反馈提交 `5b844d464fa831a33332feacffe87a6de64cc7be`。
- 当前返修基准 PPT：`{BASE}`
- 输出 PPT：`{OUT}`
- 原始 PDF：`{SRC}`
- 本轮实际修改页面：第 11、14、16、17、19、22、24 页。
- 冻结且确认未修改页面：第 8、9、10、12、13、15、18、20、21、23 页。
- 同时确认未修改页面：第 1–7、25–52 页。
- 变更页 XML 检查：仅 {changed} 与 ROUND1 基准 PPT 不同。

## 逐页修复摘要

- 第 11 页：保留顶部完整反应式，重新裁取照片区域以去除底部半截重复公式。
- 第 14 页：删除页面底部额外网址；页面可见区域只保留原课件已有来源。
- 第 16 页：保留四类键合中文解释，并恢复原课件键合科学配图区域。
- 第 17 页：标题改为 `原子轨道的电子占据`；主表格、红箭头和高亮保持不动。
- 第 19 页：重新裁取 Orbitron 图，去除上方橙色标题栏残片并保持底部轨道完整。
- 第 22 页：重新处理右侧科学图区域，去除英文正文碎片；左侧中文五层逻辑保持。
- 第 24 页：将 `N ↑ Si 原子形成晶体时` 修正为 `N 个 Si 原子形成晶体时`，其余内容保持。

## 强制自检记录

- 学生 PPT 页面中无 `源页`、`原页裁取`、`去除标题条`、`已清理`、`页面只保留`、`只保留一处`、`等比裁取`、`不叠加公式`、`真实设备照片`、`已完成中文化`、`本页已重建`、`保留原图`、`施工`、`待替换`、`占位`、`Placeholder`、`此处放图`。
- 不使用红框中文贴纸。
- 图片按 contain 方式等比放置，不做非等比拉伸。
- PPT SHA-256（渲染前）：`{h}`
''', encoding='utf-8')
print(f'round2 stage2 repair built {OUT}')
