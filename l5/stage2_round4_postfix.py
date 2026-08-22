from pathlib import Path
import subprocess, zipfile, hashlib
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT=Path('l5')
BASE=ROOT/'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND3_第16_17_24页.pptx'
SRC=ROOT/'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT=ROOT/'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND4_第16页.pptx'
EV=ROOT/'stage2_visual_review_round4'
RENDER=EV/'rendered'
COMP=EV/'comparison'
REPORT=ROOT/'BUILD_REPORT_STAGE2_ROUND4.md'

def slide_xmls(pptx):
    d={}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1,53):
            d[i]=z.read(f'ppt/slides/slide{i}.xml')
    return d

def sha256(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()

def cover(slide,x,y,w,h,color):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=color; sh.line.fill.background(); return sh

def label(slide,x,y,w,h,text,size=9,bold=False):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.background(); sh.line.fill.background()
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True; tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(1)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=text; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=RGBColor(28,28,30)
    return sh

prs=Presentation(OUT)
base_xml=slide_xmls(BASE)
slide=prs.slides[15]
# Localized masking only: cover residual fragments from original figure labels without moving or redrawing scientific figures.
# Far-right cyan remnant from the original covalent-bond caption.
cover(slide,7.48,2.35,0.92,0.28,RGBColor(255,255,255))
# Orange original metallic-bond caption remnant. Use a light gray matching the figure background and put the Chinese term in the same region.
cover(slide,6.62,6.14,2.42,0.30,RGBColor(232,232,232))
label(slide,6.88,6.12,1.78,0.34,'离域电子海',9.5,True)
# Ensure the required Chinese captions remain visible if a renderer shifts text slightly.
label(slide,5.27,2.20,2.35,0.34,'每个键由两个电子共享',8.2,False)

prs.save(OUT)
post_xml=slide_xmls(OUT)
changed=[i for i in range(1,53) if post_xml[i]!=base_xml[i]]
if changed != [16]:
    raise AssertionError(f'Only page 16 may differ from ROUND3; actual changed pages: {changed}')

# Regenerate page-16 evidence after the postfix cleanup.
pdf_dir=EV/'new_pdf'; pdf_dir.mkdir(parents=True,exist_ok=True)
subprocess.run(['libreoffice','--headless','--convert-to','pdf','--outdir',str(pdf_dir),str(OUT)],check=True)
pdf_path=pdf_dir/(OUT.stem+'.pdf')
newdoc=fitz.open(pdf_path)
doc=fitz.open(SRC)
assert newdoc.page_count==52 and doc.page_count==52
pix=newdoc[15].get_pixmap(matrix=fitz.Matrix(2.2,2.2),alpha=False)
out_png=RENDER/'page_16.png'; pix.save(out_png)
opix=doc[15].get_pixmap(matrix=fitz.Matrix(2.0,2.0),alpha=False)
orig=Image.frombytes('RGB',[opix.width,opix.height],opix.samples)
new=Image.open(out_png).convert('RGB')
h=max(orig.height,new.height)+70; w=orig.width+new.width+60
canvas=Image.new('RGB',(w,h),'white'); d=ImageDraw.Draw(canvas)
d.text((10,10),'Original PDF page 16',fill=(0,0,0)); d.text((orig.width+50,10),'New Stage 2 ROUND4 page 16',fill=(0,0,0))
canvas.paste(orig,(10,50)); canvas.paste(new,(orig.width+50,50))
comp=COMP/'page_16_original_pdf_vs_new_round4.jpg'; canvas.save(comp,quality=92)
contact=EV/'contact_sheet_stage2_round4_page_16.jpg'
thumb=Image.open(out_png).convert('RGB'); thumb.thumbnail((620,460))
cs=Image.new('RGB',(680,520),'white'); d=ImageDraw.Draw(cs); d.text((10,10),'page 16',fill=(0,0,0)); cs.paste(thumb,((680-thumb.width)//2,45)); cs.save(contact,quality=92)

hash_after=sha256(OUT)
append=f'''\n\n## ROUND4 postfix visual cleanup\n\n- 已在第 16 页局部覆盖右上晶格图中残留的英文碎片。\n- 已在第 16 页局部覆盖金属键图中残留的 `Swarm of delocalised electrons`，并在同一区域加入 `离域电子海`。\n- 冻结页复核：与 ROUND3 基准相比，仅第 16 页 XML 发生变化。\n- PPT SHA-256（postfix 后）：`{hash_after}`\n'''
REPORT.write_text(REPORT.read_text(encoding='utf-8')+append,encoding='utf-8')
subprocess.run(['git','add',str(OUT),str(out_png),str(comp),str(contact),str(REPORT)],check=True)
print('Stage 2 ROUND4 postfix cleanup applied and evidence regenerated.')
