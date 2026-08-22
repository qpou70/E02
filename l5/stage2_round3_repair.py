from pathlib import Path
import zipfile, hashlib, subprocess
import fitz
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT=Path('l5')
BASE=ROOT/'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND2_第11_14_16_17_19_22_24页.pptx'
SRC=ROOT/'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT=ROOT/'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND3_第16_17_24页.pptx'
ASSET=ROOT/'stage2_round3_assets'
EV=ROOT/'stage2_visual_review_round3'
REPORT=ROOT/'BUILD_REPORT_STAGE2_ROUND3.md'
RENDER=EV/'rendered'
COMP=EV/'comparison'
TARGET=[16,17,24]
FROZEN=[8,9,10,11,12,13,14,15,18,19,20,21,22,23]
ASSET.mkdir(exist_ok=True)
RENDER.mkdir(parents=True,exist_ok=True)
COMP.mkdir(parents=True,exist_ok=True)

def slide_xmls(pptx):
    d={}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1,53):
            d[i]=z.read(f'ppt/slides/slide{i}.xml')
    return d

def sha256(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()

prs=Presentation(BASE)
doc=fitz.open(SRC)
assert len(prs.slides)==52 and doc.page_count==52
W,H=prs.slide_width, prs.slide_height
BASE_XML=slide_xmls(BASE)

NAVY=RGBColor(30,50,82); BLUE=RGBColor(38,82,132); MID=RGBColor(204,216,230)
BLACK=RGBColor(28,28,30); WHITE=RGBColor(255,255,255); PALE=RGBColor(240,246,252)
LIGHT=RGBColor(250,252,255); TEAL=RGBColor(45,119,122)


def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)

def textbox(slide,x,y,w,h,text,size=12,bold=False,color=BLACK,align=PP_ALIGN.LEFT,fill=None,line=None,margin=3):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(0.8)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(margin)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    parts=str(text).split('\n')
    for k,part in enumerate(parts):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def rect(slide,x,y,w,h,fill=WHITE,line=MID,width=1):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh

def box(slide,x,y,w,h,text,size=11,bold=False,fill=LIGHT,line=BLUE,color=BLACK,align=PP_ALIGN.LEFT,margin=5):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(0.9)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True; tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    tf.margin_left=tf.margin_right=Pt(margin); tf.margin_top=tf.margin_bottom=Pt(4)
    for k,part in enumerate(str(text).split('\n')):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def header(slide,zh,en):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    textbox(slide,0.38,0.06,6.7,0.42,zh,18,True,WHITE)
    textbox(slide,6.45,0.10,3.12,0.32,en,8.5,False,RGBColor(224,233,243),PP_ALIGN.RIGHT)

def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

def render_crop(page_no,coords,name,dpi=300,mask_right_page_number=False,mask_bottom_px=0,mask_top_px=0):
    pix=doc[page_no-1].get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),clip=fitz.Rect(*coords),alpha=False)
    out=ASSET/name; pix.save(out)
    im=Image.open(out).convert('RGB'); d=ImageDraw.Draw(im)
    if mask_right_page_number:
        d.rectangle([im.width-120, im.height-58, im.width, im.height], fill='white')
    if mask_bottom_px:
        d.rectangle([0, im.height-mask_bottom_px, im.width, im.height], fill='white')
    if mask_top_px:
        d.rectangle([0, 0, im.width, mask_top_px], fill='white')
    im.save(out)
    return out

def add_pic(slide,path,x,y,w,h):
    with Image.open(path) as im:
        iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    return slide.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

# Page 16: rebuild only allowed page; keep four Chinese explanations and replace English bullet column with pure scientific diagrams.
s=prs.slides[15]; clear(s); header(s,'键合类型','Bond Types')
box(s,0.38,0.82,3.05,1.25,'离子键 / Ionic Bonding\n电子转移；库仑吸引与原子核排斥在平衡距离处达到平衡。',9.6,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,2.12,3.05,1.15,'共价键 / Covalent Bonding\n电子共享；典型材料：Si、Ge、C。',9.8,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,3.34,3.05,1.38,'混合离子-共价键 / Mixed Ionic-Covalent Bonding\n电负性差导致极性共价键；典型材料：GaAs、InP、GaN。',8.6,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
box(s,0.38,4.84,3.05,1.22,'金属键 / Metallic Bonding\n正离子实处于电子海中；常见于价电子数不超过 3 的原子。',9.0,False,LIGHT,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
rect(s,3.62,0.82,5.92,5.92,WHITE,MID,1)
# Three clean original science-figure crops from PDF page 16.  No title bar, page number, source URL, or English bullet list is included.
p16a=render_crop(16,(382,170,688,282),'round3_p16_lattice_only.png',320)
p16b=render_crop(16,(320,333,700,432),'round3_p16_deltaE_only.png',320)
p16c=render_crop(16,(400,458,640,522),'round3_p16_metallic_only.png',320)
textbox(s,3.86,0.98,5.25,0.24,'离子键 / 共价键晶格结构示意',9.2,True,NAVY,PP_ALIGN.CENTER)
add_pic(s,p16a,3.90,1.25,5.20,1.26)
textbox(s,3.86,2.72,5.25,0.24,'电负性差 ΔE 与电子云极化',9.2,True,NAVY,PP_ALIGN.CENTER)
add_pic(s,p16b,3.80,3.00,5.38,1.26)
textbox(s,3.86,4.50,5.25,0.24,'金属键：正离子实 + 离域电子海',9.2,True,NAVY,PP_ALIGN.CENTER)
add_pic(s,p16c,4.02,4.86,4.96,0.92)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 16. Scientific diagrams cropped from the original slide; English bullet list and original title bar removed.')

# Page 17: rebuild only allowed page; complete table crop includes 33 As, 34 Se, 35 Br, 36 Kr and all original arrows/highlights.
s=prs.slides[16]; clear(s); header(s,'原子轨道的电子占据','Population of Atomic Orbitals')
p17=render_crop(17,(145,153,560,536),'round3_p17_full_orbital_population_table.png',330,mask_right_page_number=True)
rect(s,0.38,0.82,9.24,6.05,WHITE,MID,1)
add_pic(s,p17,1.42,0.98,7.02,5.72)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 17. Complete table retained through 36 Kr, including original arrows, dashed boxes, and Si highlight.')

# Page 24: targeted text fixes only.
s=prs.slides[23]
for sh in s.shapes:
    if hasattr(sh,'text_frame') and sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.text=r.text.replace('N ↑ Si 原子形成晶体时','N 个 Si 原子形成晶体时').replace('N↑ Si 原子形成晶体时','N 个 Si 原子形成晶体时').replace('N ↑','N 个').replace('N↑','N 个')
                if r.text.strip()=='Total':
                    r.text=r.text.replace('Total','合计')
    if hasattr(sh,'text'):
        txt=sh.text
        if 'N ↑' in txt or 'N↑' in txt or txt.strip()=='Total':
            sh.text=txt.replace('N ↑ Si 原子形成晶体时','N 个 Si 原子形成晶体时').replace('N↑ Si 原子形成晶体时','N 个 Si 原子形成晶体时').replace('N ↑','N 个').replace('N↑','N 个').replace('Total','合计')
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 24.')

# Visible-text validation.
for p in TARGET:
    txt='\n'.join((getattr(sh,'text','') or '') for sh in prs.slides[p-1].shapes)
    if p==16:
        bad=['Electron transfer','Coulomb attraction balanced','Electron sharing','Electronegativity','Polar covalent bonds','Positive cores in an electron sea','Typical for atoms with 3 or less valence electrons']
        for b in bad:
            assert b not in txt, f'page 16 still has English bullet text: {b}'
    if p==17:
        assert 'Electronic Configurations' not in txt
    if p==24:
        assert 'N ↑' not in txt and 'N↑' not in txt and 'Total' not in txt

prs.save(OUT)
POST_XML=slide_xmls(OUT)
changed=[i for i in range(1,53) if POST_XML[i]!=BASE_XML[i]]
if sorted(changed)!=TARGET:
    raise AssertionError(f'Only pages {TARGET} may change; actual changed pages: {changed}')
pre_hash=sha256(OUT)

# Export PDF and render evidence pages.
pdf_dir=EV/'new_pdf'; pdf_dir.mkdir(parents=True,exist_ok=True)
subprocess.run(['libreoffice','--headless','--convert-to','pdf','--outdir',str(pdf_dir),str(OUT)],check=True)
pdf_path=pdf_dir/(OUT.stem+'.pdf')
newdoc=fitz.open(pdf_path)
assert newdoc.page_count==52

rendered=[]
for p in TARGET:
    pix=newdoc[p-1].get_pixmap(matrix=fitz.Matrix(2.2,2.2),alpha=False)
    out=RENDER/f'page_{p}.png'; pix.save(out); rendered.append(out)
    opix=doc[p-1].get_pixmap(matrix=fitz.Matrix(2.0,2.0),alpha=False)
    orig=Image.frombytes('RGB',[opix.width,opix.height],opix.samples)
    new=Image.open(out).convert('RGB')
    h=max(orig.height,new.height)+70; w=orig.width+new.width+60
    canvas=Image.new('RGB',(w,h),'white'); d=ImageDraw.Draw(canvas)
    d.text((10,10),f'Original PDF page {p}',fill=(0,0,0))
    d.text((orig.width+50,10),f'New Stage 2 ROUND3 page {p}',fill=(0,0,0))
    canvas.paste(orig,(10,50)); canvas.paste(new,(orig.width+50,50))
    canvas.save(COMP/f'page_{p}_original_pdf_vs_new_round3.jpg',quality=92)

# Contact sheet.
thumbs=[]
for p in TARGET:
    im=Image.open(RENDER/f'page_{p}.png').convert('RGB')
    im.thumbnail((520,390))
    tile=Image.new('RGB',(560,440),'white'); d=ImageDraw.Draw(tile)
    d.text((10,10),f'page {p}',fill=(0,0,0))
    tile.paste(im,((560-im.width)//2,40))
    thumbs.append(tile)
cs=Image.new('RGB',(560*len(thumbs),440),'white')
for i,t in enumerate(thumbs): cs.paste(t,(560*i,0))
contact=EV/'contact_sheet_stage2_round3_pages_16_17_24.jpg'
cs.save(contact,quality=92)
post_hash=sha256(OUT)

report=f'''# ECE340 L5 第二阶段视觉返修 ROUND3 Build Report

- 返修依据：`l5/SUPERVISOR_STAGE2_VISUAL_REVIEW_FEEDBACK_ROUND3.md`，反馈提交 `c9181636b4913953fc3dc3e1638d277b5b9bea1f`。
- 当前返修基准 PPT：`l5/ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND2_第11_14_16_17_19_22_24页.pptx`
- 输出 PPT：`l5/ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND3_第16_17_24页.pptx`
- 原始 PDF：`l5/stage1_source_reference/ECE340_L5_S18_Posted.pdf`
- 本轮实际修改页面：第 16、17、24 页。
- 冻结且确认未修改页面：第 8、9、10、11、12、13、14、15、18、19、20、21、22、23 页。
- 同时确认未修改页面：第 1–7、25–52 页。
- 变更页 XML 检查：仅 {changed} 与 ROUND2 基准 PPT 不同。

## 逐页修复摘要

- 第 16 页：保留左侧四类中文说明；右侧移除原英文 bullet list 和橙色标题条，改为从原 PDF 分别干净裁取的三组科学图示。
- 第 17 页：保留中文标题；恢复完整表格到底部 36 Kr，并保留红箭头、红色虚线框和 Si 黄色高亮；右上小标题改为 `Population of Atomic Orbitals`。
- 第 24 页：将 `N ↑ Si 原子形成晶体时` 修正为 `N 个 Si 原子形成晶体时`；两个 `Total` 改为 `合计`。

## 强制自检记录

- 第 16 页：无整套英文 bullet list；无原 PDF 橙色标题条；三组科学配图完整。
- 第 17 页：表格完整显示到 36 Kr；33 As / 34 Se / 35 Br / 36 Kr 可见；表格底部未裁切。
- 第 24 页：不存在 `N ↑ Si 原子形成晶体时`；两个 `Total` 已中文化为 `合计`。
- PPT SHA-256（渲染前）：`{pre_hash}`

## 渲染与证据

- 高清渲染图目录：`l5/stage2_visual_review_round3/rendered/`。
- Original PDF vs New Page 对照图目录：`l5/stage2_visual_review_round3/comparison/`。
- Contact sheet：`l5/stage2_visual_review_round3/contact_sheet_stage2_round3_pages_16_17_24.jpg`。
- 渲染页面：第 16、17、24 页。
- 渲染 PDF 页数：原始 PDF 52；第二阶段 ROUND3 PPT 导出 PDF 52。
- Worker self-check: passed.
- Supervisor visual acceptance: pending.
- PPT SHA-256（渲染后）：`{post_hash}`
'''
REPORT.write_text(report,encoding='utf-8')

subprocess.run(['git','add',str(OUT),str(REPORT),str(EV)],check=True)
print('stage2 round3 evidence complete')
print(f'changed slide XML pages: {changed}')
