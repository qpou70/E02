from pathlib import Path
import re, zipfile, hashlib
import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

ROOT = Path("l5")
BASE = ROOT / "ECE340_L5_S18_Posted_中文忠实重建_视觉返修版_第8-24页.pptx"
SRC = ROOT / "stage1_source_reference/ECE340_L5_S18_Posted.pdf"
OUT = ROOT / "ECE340_L5_S18_Posted_中文忠实重建_第一阶段严格返修_第9_12_18_20_23页.pptx"
ASSET = ROOT / "stage1_strict_assets"
REPORT = ROOT / "BUILD_REPORT_08_24_STAGE1_STRICT.md"
TARGET = [9, 12, 18, 20, 23]
ASSET.mkdir(exist_ok=True)

assert BASE.exists(), BASE
assert SRC.exists(), SRC
prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52
assert doc.page_count == 52
W, H = prs.slide_width, prs.slide_height

NAVY = RGBColor(30, 50, 82)
BLUE = RGBColor(38, 82, 132)
TEAL = RGBColor(45, 119, 122)
GOLD = RGBColor(208, 153, 55)
GRAY = RGBColor(92, 98, 108)
BLACK = RGBColor(28, 28, 30)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(240, 246, 252)
PALETEAL = RGBColor(237, 248, 247)
LIGHT = RGBColor(250, 252, 255)
MID = RGBColor(204, 216, 230)

def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)

def textbox(slide, x, y, w, h, text, size=14, bold=False, color=BLACK,
            align=PP_ALIGN.LEFT, fill=None, line=None, margin=4):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(margin)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Noto Sans CJK SC"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh

def box(slide, x, y, w, h, text="", size=13, bold=False, fill=LIGHT, line=BLUE,
        color=BLACK, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1.05)
    tf = sh.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(7)
    tf.margin_top = tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = "Noto Sans CJK SC"; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sh

def rect(slide, x, y, w, h, fill=WHITE, line=MID, width=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(width)
    return sh

def oval(slide, x, y, w, h, fill=PALE, line=BLUE, width=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(width)
    return sh

def connector(slide, x1, y1, x2, y2, color=GRAY, width=1.2):
    sh = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    sh.line.color.rgb = color; sh.line.width = Pt(width)
    return sh

def right_arrow(slide, x, y, w, h, fill=TEAL):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def header(slide, zh, en):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.60))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    textbox(slide, 0.48, 0.06, 8.7, 0.44, zh, 21, True, WHITE)
    textbox(slide, 8.55, 0.10, 4.20, 0.34, en, 10, False, RGBColor(224,233,243), PP_ALIGN.RIGHT)

def render_crop(page_no, rect_points, name, dpi=230):
    page = doc[page_no - 1]
    pix = page.get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72), clip=fitz.Rect(*rect_points), alpha=False)
    out = ASSET / name
    pix.save(out)
    return out

def add_picture_contain(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(w/iw, h/ih)
    pw, ph = iw*ratio, ih*ratio
    px, py = x+(w-pw)/2, y+(h-ph)/2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# PAGE 9
s = prs.slides[8]; clear(s)
header(s, "硅的气相外延（VPE）", "Silicon Vapor Phase Epitaxy")
box(s, 0.55, 0.76, 5.70, 0.50, "SiCl₄ + 2H₂ ⇌ Si + 4HCl", 17, False, WHITE, BLUE, BLACK, PP_ALIGN.CENTER)
box(s, 6.45, 0.76, 4.90, 0.50, "SiH₄ → Si + 2H₂", 17, False, WHITE, BLUE, BLACK, PP_ALIGN.CENTER)
photo = render_crop(9, (36,126,225,288), "p09_equipment.png")
rect(s, 0.55, 1.42, 3.45, 4.70, WHITE, MID, 1.1)
add_picture_contain(s, photo, 0.66, 1.54, 3.23, 4.08)
textbox(s, 0.76, 5.68, 3.02, 0.30, "VPE 外延设备", 11, True, NAVY, PP_ALIGN.CENTER)
rect(s, 4.22, 1.42, 8.55, 4.70, RGBColor(248,251,255), BLUE, 1.2)
textbox(s, 7.12, 1.53, 2.25, 0.30, "石英反应腔", 12, True, NAVY, PP_ALIGN.CENTER)
textbox(s, 4.40, 1.83, 1.05, 0.27, "气体入口", 10, True, TEAL, PP_ALIGN.CENTER)
right_arrow(s, 4.48, 2.13, 1.18, 0.24, TEAL)
for xx in (5.72, 6.62, 7.52, 8.42, 9.32, 10.22):
    right_arrow(s, xx, 2.16, 0.70, 0.18, TEAL)
rect(s, 5.76, 1.90, 0.14, 1.18, RGBColor(236,224,191), GOLD, 1)
textbox(s, 5.28, 3.02, 1.12, 0.26, "挡流板", 9, True, GRAY, PP_ALIGN.CENTER)
connector(s, 5.77, 3.04, 5.83, 2.58, GRAY, 0.9)
rect(s, 7.05, 3.75, 3.80, 0.23, RGBColor(112,121,134), GRAY, 0.9)
for xx in (7.22, 8.10, 8.98, 9.86):
    oval(s, xx, 3.36, 0.62, 0.22, RGBColor(193,224,246), BLUE, 0.8)
textbox(s, 8.13, 3.00, 1.70, 0.29, "硅片 / 晶圆", 10, True, NAVY, PP_ALIGN.CENTER)
connector(s, 9.05, 3.28, 9.28, 3.45, NAVY, 0.9)
textbox(s, 8.11, 4.03, 1.76, 0.26, "基座（Susceptor）", 9, True, GRAY, PP_ALIGN.CENTER)
rect(s, 8.66, 3.98, 0.62, 0.92, RGBColor(204,210,216), GRAY, 0.9)
textbox(s, 9.48, 4.35, 1.65, 0.27, "支座 / Pedestal", 9, True, GRAY)
connector(s, 9.50, 4.48, 9.27, 4.44, GRAY, 0.9)
for i in range(4):
    right_arrow(s, 6.93+i*0.72, 5.05, 0.55, 0.17, GOLD)
textbox(s, 7.12, 5.34, 1.68, 0.28, "RF 加热", 10, True, GOLD, PP_ALIGN.CENTER)
right_arrow(s, 10.93, 2.13, 1.35, 0.24, TEAL)
textbox(s, 11.12, 1.82, 1.15, 0.28, "排气口", 10, True, TEAL, PP_ALIGN.CENTER)
textbox(s, 4.58, 5.72, 7.55, 0.30, "气流沿反应腔流过加热的晶圆表面，再由排气口排出。", 11, False, BLACK, PP_ALIGN.CENTER)
notes(s, "VPE 反应气体沿反应腔通过加热晶圆表面。两条反应式严格按原页保留。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 9.")

# PAGE 12
s = prs.slides[11]; clear(s)
header(s, "分子束外延（MBE）", "Molecular Beam Epitaxy")
rect(s, 0.55, 0.88, 6.00, 5.96, RGBColor(248,251,255), BLUE, 1.2)
textbox(s, 2.12, 1.02, 2.92, 0.30, "超高真空生长腔", 13, True, NAVY, PP_ALIGN.CENTER)
rect(s, 2.66, 1.68, 1.92, 0.26, RGBColor(183,216,243), BLUE, 0.9)
textbox(s, 2.65, 1.99, 1.95, 0.27, "GaAs 衬底", 10, True, NAVY, PP_ALIGN.CENTER)
rect(s, 3.18, 1.39, 0.86, 0.16, RGBColor(188,192,198), GRAY, 0.8)
textbox(s, 4.18, 1.34, 1.18, 0.26, "衬底支架", 9, False, GRAY)
connector(s, 4.20, 1.47, 4.02, 1.52, GRAY, 0.8)
sources = [("Si",0.90),("Al",1.82),("Ga",2.74),("As",3.66),("Be",4.58)]
for lab, x in sources:
    rect(s, x, 5.66, 0.52, 0.50, RGBColor(239,241,244), GRAY, 0.9)
    textbox(s, x, 5.72, 0.52, 0.22, lab, 9, True, NAVY, PP_ALIGN.CENTER)
    rect(s, x-0.02, 5.22, 0.58, 0.07, RGBColor(150,156,164), GRAY, 0.3)
    connector(s, x+0.26, 5.20, 3.62, 2.24, GOLD, 1.2)
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.51), Inches(2.18), Inches(0.16), Inches(0.14))
    tri.fill.solid(); tri.fill.fore_color.rgb = GOLD; tri.line.fill.background()
textbox(s, 0.86, 6.22, 4.42, 0.28, "束源：Si / Al / Ga / As / Be", 9, True, GRAY, PP_ALIGN.CENTER)
textbox(s, 4.82, 4.90, 1.20, 0.28, "快门", 9, True, GRAY, PP_ALIGN.CENTER)
connector(s, 4.88, 5.09, 4.57, 5.25, GRAY, 0.8)
textbox(s, 1.23, 3.31, 1.42, 0.28, "分子束", 10, True, GOLD, PP_ALIGN.CENTER)
connector(s, 2.42, 3.44, 3.00, 3.03, GOLD, 0.9)
micro = render_crop(12, (390,126,720,535), "p12_right_source.png")
rect(s, 6.75, 0.88, 6.00, 5.96, WHITE, MID, 1.1)
add_picture_contain(s, micro, 6.90, 1.02, 5.70, 5.66)
notes(s, "左图保留 Si / Al / Ga / As / Be 束源、快门、分子束与 GaAs 衬底的对应关系。右图直接使用原页真实图像裁取，不添加中文贴纸；图内 4×4 / GaAs substrate / 10 nm / <100> 按原图保留。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 12.")

# PAGE 18
s = prs.slides[17]; clear(s)
header(s, "硅的芯层电子与价电子", "Silicon Core and Valence Electrons")
p18 = render_crop(18, (54,126,738,535), "p18_complete_source_body.png", 250)
rect(s, 0.48, 0.82, 12.38, 6.20, WHITE, MID, 1.0)
add_picture_contain(s, p18, 0.62, 0.96, 12.10, 5.92)
notes(s, "本页两幅高信息密度原图完整保留，不在英文原图上叠加中文框。讲解重点：+14；1s / 2s / 2p / 3s / 3p；芯层电子、价电子、价轨道、首次激发轨道、价能级与电离/零能级的关系。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 18.")

# PAGE 20
s = prs.slides[19]; clear(s)
header(s, "sp³ 杂化与四面体几何", "sp³ Hybridization")
p20 = render_crop(20, (42,126,754,535), "p20_formula_geometry_source.png", 260)
rect(s, 0.48, 0.82, 12.38, 6.20, WHITE, MID, 1.0)
add_picture_contain(s, p20, 0.62, 0.96, 12.10, 5.92)
notes(s, "原页四条 sp³ 杂化公式与 109.5° 四面体几何直接使用源页裁图保留，不重新誊写、不改系数。原页还指出应变会扭曲键角，氨等分子的键角略小。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 20.")

# PAGE 23
s = prs.slides[22]; clear(s)
header(s, "氢分子的成键与反键轨道", "Bonding and Antibonding Orbitals in Hydrogen")
box(s, 0.50, 0.86, 3.15, 2.15,
    "元素氢：1s¹\nH #1: 2 states, 1 electron\nH #2: 2 states, 1 electron\nH₂: 4 states, 2 electrons",
    12, False, PALE, BLUE, BLACK)
box(s, 0.50, 3.27, 3.15, 1.12,
    "两个 H 1s 原子轨道组合后形成两个分子轨道：\n低能成键轨道 + 高能反键轨道。",
    11, True, WHITE, BLUE, BLACK, PP_ALIGN.CENTER)
box(s, 0.50, 4.64, 3.15, 1.56,
    "H₂ 的 2 个电子在基态优先占据低能成键轨道，因此形成稳定 H–H 键。",
    11, False, PALETEAL, TEAL, BLACK, PP_ALIGN.CENTER)
rect(s, 3.92, 0.86, 5.42, 5.96, RGBColor(251,253,255), MID, 1.0)
textbox(s, 5.30, 1.00, 2.70, 0.30, "分子轨道能级关系", 13, True, NAVY, PP_ALIGN.CENTER)
connector(s, 4.32, 3.78, 5.32, 3.78, GRAY, 1.8)
connector(s, 7.96, 3.78, 8.96, 3.78, GRAY, 1.8)
textbox(s, 4.28, 3.38, 1.08, 0.28, "H #1 1s", 9, True, GRAY, PP_ALIGN.CENTER)
textbox(s, 7.92, 3.38, 1.08, 0.28, "H #2 1s", 9, True, GRAY, PP_ALIGN.CENTER)
connector(s, 5.80, 2.08, 7.48, 2.08, NAVY, 2.2)
textbox(s, 5.62, 1.64, 2.06, 0.30, "反键 σ*1s", 10, True, NAVY, PP_ALIGN.CENTER)
connector(s, 5.80, 5.17, 7.48, 5.17, TEAL, 2.2)
textbox(s, 5.62, 5.34, 2.06, 0.30, "成键 σ1s", 10, True, TEAL, PP_ALIGN.CENTER)
for a,b,c,d in [(5.32,3.78,5.80,2.08),(5.32,3.78,5.80,5.17),(7.96,3.78,7.48,2.08),(7.96,3.78,7.48,5.17)]:
    connector(s,a,b,c,d,GRAY,1.0)
textbox(s, 6.24, 4.69, 0.38, 0.34, "↑", 17, True, TEAL, PP_ALIGN.CENTER)
textbox(s, 6.65, 4.69, 0.38, 0.34, "↓", 17, True, TEAL, PP_ALIGN.CENTER)
up = s.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(8.90), Inches(1.62), Inches(0.26), Inches(3.92))
up.fill.solid(); up.fill.fore_color.rgb = NAVY; up.line.fill.background()
textbox(s, 8.32, 1.16, 0.92, 0.40, "Higher\nEnergy", 8, True, NAVY, PP_ALIGN.CENTER)
textbox(s, 8.32, 5.58, 0.92, 0.40, "Lower\nEnergy", 8, True, TEAL, PP_ALIGN.CENTER)
rect(s, 9.56, 0.86, 3.20, 5.96, WHITE, MID, 1.0)
textbox(s, 9.78, 1.00, 2.76, 0.30, "轨道与电子密度", 12, True, NAVY, PP_ALIGN.CENTER)
oval(s, 9.90, 1.67, 0.92, 0.92, RGBColor(229,239,250), BLUE, 1)
oval(s, 11.47, 1.67, 0.92, 0.92, RGBColor(247,236,216), GOLD, 1)
textbox(s, 10.20, 1.95, 0.32, 0.24, "+", 11, True, NAVY, PP_ALIGN.CENTER)
textbox(s, 11.77, 1.95, 0.32, 0.24, "−", 11, True, GOLD, PP_ALIGN.CENTER)
rect(s, 11.10, 1.55, 0.04, 1.15, GRAY, GRAY, 0)
textbox(s, 9.78, 2.70, 2.78, 0.58, "反键：两核之间出现节点\n电子不位于两核之间", 9, True, NAVY, PP_ALIGN.CENTER)
oval(s, 9.96, 4.03, 1.17, 0.94, RGBColor(220,240,238), TEAL, 1)
oval(s, 11.16, 4.03, 1.17, 0.94, RGBColor(220,240,238), TEAL, 1)
oval(s, 10.68, 4.14, 0.92, 0.72, RGBColor(195,227,224), TEAL, 0.8)
textbox(s, 9.78, 5.05, 2.78, 0.58, "成键：两核之间电子密度增加\n系统能量降低", 9, True, TEAL, PP_ALIGN.CENTER)
notes(s, "H #1 与 H #2 的 1s 原子轨道组合为低能成键轨道和高能反键轨道。成键轨道在两核之间电子密度增加；反键轨道在两核之间形成节点。H₂ 的两个电子在基态占据成键轨道。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 23.")

for p in TARGET:
    slide = prs.slides[p-1]
    for sh in slide.shapes:
        assert sh.left >= 0 and sh.top >= 0
        assert sh.left + sh.width <= W + 1000, (p, sh.name, "x overflow")
        assert sh.top + sh.height <= H + 1000, (p, sh.name, "y overflow")

prs.save(OUT)

with zipfile.ZipFile(BASE) as zb, zipfile.ZipFile(OUT) as zo:
    for p in range(1, 53):
        a = zb.read(f"ppt/slides/slide{p}.xml")
        b = zo.read(f"ppt/slides/slide{p}.xml")
        if p in TARGET:
            assert a != b, f"target page {p} unchanged"
        else:
            assert a == b, f"non-target page {p} changed"
    for p in range(1, 53):
        n = f"ppt/notesSlides/notesSlide{p}.xml"
        if p not in TARGET and n in zb.namelist() and n in zo.namelist():
            assert zb.read(n) == zo.read(n), f"non-target notes page {p} changed"

check = Presentation(OUT)
banned = [
    "原设备照片","右侧显微图","原图示意","待替换图片","此处保留原图",
    "本页已重建","已删除残留","保留原始信息","为避免 PowerPoint 压字",
    "本页采用中文示意","右图保留","已完成中文化","已重做","已改为"
]
for p in TARGET:
    slide = check.slides[p-1]
    texts = []
    for sh in slide.shapes:
        assert sh.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER, (p, sh.name, "placeholder")
        if getattr(sh, "has_text_frame", False):
            texts.append(sh.text)
    joined = "\n".join(texts)
    for term in banned:
        assert term not in joined, (p, term)

t9 = "\n".join(sh.text for sh in check.slides[8].shapes if getattr(sh, "has_text_frame", False))
assert "SiCl₄ + 2H₂ ⇌ Si + 4HCl" in t9
assert "SiH₄ → Si + 2H₂" in t9
for term in ["石英反应腔","气体入口","挡流板","硅片 / 晶圆","基座（Susceptor）","支座 / Pedestal","RF 加热","排气口"]:
    assert term in t9, term

t12 = "\n".join(sh.text for sh in check.slides[11].shapes if getattr(sh, "has_text_frame", False))
for term in ["Si / Al / Ga / As / Be","GaAs 衬底","快门","分子束"]:
    assert term in t12, term

t23 = "\n".join(sh.text for sh in check.slides[22].shapes if getattr(sh, "has_text_frame", False))
for term in ["H #1: 2 states, 1 electron","H₂: 4 states, 2 electrons","Higher","Lower","成键","反键"]:
    assert term in t23, term

src20 = doc[19].get_text("text").replace("\n"," ")
for seq in [
    "ψ 1 = 1", "ψ s +ψ px +ψ py +ψ pz",
    "ψ 2 = 1", "ψ s −ψ px −ψ py +ψ pz",
    "ψ 3 = 1", "ψ s +ψ px −ψ py −ψ pz",
    "ψ 4 = 1", "ψ s −ψ px +ψ py −ψ pz",
    "109.5° Bond"
]:
    assert seq in src20, seq

digest = hashlib.sha256(OUT.read_bytes()).hexdigest()
REPORT.write_text(
    "# ECE340 L5 第一阶段严格视觉返修 Build Report\n\n"
    f"- 基准 PPT：`{BASE.as_posix()}`\n"
    f"- 原始讲义 PDF：`{SRC.as_posix()}`\n"
    f"- 阶段输出 PPT：`{OUT.as_posix()}`\n"
    "- 实际修改页：9、12、18、20、23。\n"
    "- 第 8–24 页中保持未改：8、10、11、13、14、15、16、17、19、21、22、24。\n"
    "- 其他保持未改：1–7、25–52。\n"
    "- 页 9：保留原页真实设备照片；重建 VPE 反应器对象关系；两条反应式逐字断言。\n"
    "- 页 12：重建 Si/Al/Ga/As/Be 束源—快门—分子束—衬底关系；右侧使用原页真实裁图。\n"
    "- 页 18：两幅高密度关系图使用原页正文干净裁图完整保留；中文讲解置备注。\n"
    "- 页 20：四条 sp³ 公式与 109.5° 直接使用原页正文裁图，不重新誊写或换公式。\n"
    "- 页 23：重建状态数、原子 1s—分子轨道能级、成键/反键、Higher/Lower Energy 与电子密度对应。\n"
    "- Placeholder：无；目标页已检查 shape_type。\n"
    "- 红框中文贴纸：无。\n"
    "- 施工说明：无；已执行禁词扫描。\n"
    "- 科学内容/公式：未修改；页 20 原始公式文本已从源 PDF 逐项校验，页面直接使用源裁图。\n"
    "- 非目标页 slide XML：与基准逐页 byte-for-byte 相同。\n"
    f"- PPT SHA-256（渲染前）：`{digest}`\n",
    encoding="utf-8"
)
print(OUT)
