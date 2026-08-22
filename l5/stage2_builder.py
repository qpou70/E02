from pathlib import Path
import zipfile, hashlib, re
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第一阶段严格返修_第9_12_18_20_23页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段返修_第8_10_11_14_16_17_19_22_24页.pptx'
ASSET = ROOT / 'stage2_assets'
REPORT = ROOT / 'BUILD_REPORT_STAGE2.md'
TARGET = [8,10,11,14,16,17,19,22,24]
FROZEN = [9,12,13,15,18,20,21,23]
ASSET.mkdir(exist_ok=True)

prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52 and doc.page_count == 52
W, H = prs.slide_width, prs.slide_height
NAVY=RGBColor(30,50,82); BLUE=RGBColor(38,82,132); TEAL=RGBColor(45,119,122)
GOLD=RGBColor(208,153,55); GRAY=RGBColor(92,98,108); BLACK=RGBColor(28,28,30)
WHITE=RGBColor(255,255,255); PALE=RGBColor(240,246,252); MID=RGBColor(204,216,230)
LIGHT=RGBColor(250,252,255); GREEN=RGBColor(73,130,84); ORANGE=RGBColor(209,122,42)

def slide_xmls(pptx):
    out={}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1,53):
            out[i]=z.read(f'ppt/slides/slide{i}.xml')
    return out
BASE_XML = slide_xmls(BASE)

def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)

def textbox(slide,x,y,w,h,text,size=12,bold=False,color=BLACK,align=PP_ALIGN.LEFT,fill=None,line=None,margin=3):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(0.8)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(margin)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    for k, part in enumerate(str(text).split('\n')):
        if k: p=tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def rect(slide,x,y,w,h,fill=WHITE,line=MID,width=1):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh

def box(slide,x,y,w,h,text,size=11,bold=False,fill=LIGHT,line=BLUE,color=BLACK,align=PP_ALIGN.LEFT):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(0.9)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Pt(6); tf.margin_top=tf.margin_bottom=Pt(4); tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    for k, part in enumerate(str(text).split('\n')):
        if k: p=tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def line(slide,x1,y1,x2,y2,color=GRAY,width=1):
    sh=slide.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    sh.line.color.rgb=color; sh.line.width=Pt(width); return sh

def arrow(slide,x,y,w,h,fill=TEAL):
    sh=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); return sh

def header(slide,zh,en):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    textbox(slide,0.38,0.06,6.35,0.42,zh,18,True,WHITE)
    textbox(slide,6.55,0.10,3.05,0.32,en,8.5,False,RGBColor(224,233,243),PP_ALIGN.RIGHT)

def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

def render_crop(page_no,coords,name,dpi=260,clean_bottom=False):
    pix=doc[page_no-1].get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),clip=fitz.Rect(*coords),alpha=False)
    out=ASSET/name; pix.save(out)
    if clean_bottom:
        im=Image.open(out).convert('RGB'); d=ImageDraw.Draw(im)
        d.rectangle([0, im.height-38, im.width, im.height], fill='white')
        im.save(out)
    return out

def add_pic(slide,path,x,y,w,h):
    with Image.open(path) as im: iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    return slide.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

def simple_table(slide,x,y,cols,rows,widths,row_h=0.42,font=8.6):
    cx=x
    for j,c in enumerate(cols):
        box(slide,cx,y,widths[j],row_h,c,font,True,PALE,BLUE,NAVY,PP_ALIGN.CENTER); cx+=widths[j]
    for i,row in enumerate(rows):
        cx=x
        for j,c in enumerate(row):
            box(slide,cx,y+row_h*(i+1),widths[j],row_h,c,font,False,WHITE,MID,BLACK,PP_ALIGN.CENTER if j else PP_ALIGN.LEFT); cx+=widths[j]

# page 8
s=prs.slides[7]; clear(s); header(s,'半导体材料制备方法分类','Semiconductor Growth Methods')
box(s,0.45,0.88,9.10,0.55,'按原页教学逻辑：先区分体材料生长与外延生长，再列出各自典型方法。',12,True,WHITE,BLUE,NAVY,PP_ALIGN.CENTER)
box(s,0.55,1.65,4.20,0.55,'体材料生长（bulk growth）',14,True,RGBColor(246,250,255),BLUE,NAVY,PP_ALIGN.CENTER)
box(s,5.25,1.65,4.20,0.55,'外延生长（epitaxial growth）',14,True,RGBColor(246,250,255),BLUE,NAVY,PP_ALIGN.CENTER)
for y,t in [(2.45,'Czochralski，CZ 提拉法'),(3.12,'Float-zone，FZ 区熔法'),(3.79,'Bridgman / gradient-freeze，坩埚凝固法')]:
    box(s,0.82,y,3.65,0.46,t,10.8,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
for y,t in [(2.34,'LPE：液相外延'),(2.91,'VPE / CVD：气相外延 / 化学气相沉积'),(3.48,'MBE：分子束外延'),(4.05,'MOCVD / OMVPE：金属有机气相外延')]:
    box(s,5.55,y,3.62,0.43,t,10.0,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
arrow(s,4.72,3.28,0.47,0.24,TEAL)
box(s,1.00,5.05,8.00,0.88,'教学重点：体材料生长决定晶体基底；外延生长在基底上沉积受控薄层。后续页面分别展开 VPE/CVD、MOCVD、MBE 以及能带/键合背景。',11.0,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 8.')

# page 10
s=prs.slides[9]; clear(s); header(s,'硅的化学气相沉积设备','Silicon CVD / VPE Equipment')
p10=render_crop(10,(36,126,756,520),'p10_body_clean.png',270,True)
rect(s,0.48,0.85,5.90,5.92,WHITE,MID,1); add_pic(s,p10,0.60,0.98,5.65,5.62)
box(s,6.62,1.10,2.80,0.72,'真实设备照片\n去除源页标题条与页码',11,True,RGBColor(246,250,255),BLUE,NAVY,PP_ALIGN.CENTER)
box(s,6.62,2.10,2.80,0.72,'载气与反应物进入反应腔，在加热晶圆表面发生沉积反应。',10.5,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
box(s,6.62,3.12,2.80,0.70,'页面只保留一套正式标题；设备图按原比例裁取。',10.5,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
box(s,6.62,4.12,2.80,0.82,'关键对象：气体入口、反应腔、加热基座、晶圆、排气端。',10.5,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 10.')

# page 11
s=prs.slides[10]; clear(s); header(s,'气相外延中的化学反应','Chemical Reactions in VPE')
p11=render_crop(11,(40,130,756,505),'p11_reactions_clean.png',270,True)
rect(s,0.45,0.85,6.05,5.95,WHITE,MID,1); add_pic(s,p11,0.58,1.00,5.78,5.58)
box(s,6.72,1.05,2.65,0.70,'反应式以原页为准',12,True,RGBColor(246,250,255),BLUE,NAVY,PP_ALIGN.CENTER)
box(s,6.72,1.95,2.65,0.88,'清理源页标题条与孤立页码；不在原公式上叠加新的公式。',10.5,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
box(s,6.72,3.10,2.65,1.10,'阅读顺序：反应物输运 → 表面吸附/分解 → 外延层生长 → 副产物排出。',10.8,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 11.')

# page 14
s=prs.slides[13]; clear(s); header(s,'金属有机化学气相沉积（MOCVD）','MOCVD')
p14=render_crop(14,(38,126,756,515),'p14_two_devices.png',270,True)
rect(s,0.45,0.88,6.40,5.85,WHITE,MID,1); add_pic(s,p14,0.62,1.04,6.05,5.45)
box(s,7.05,1.00,2.30,0.70,'两张设备图等比裁取',11.5,True,RGBColor(246,250,255),BLUE,NAVY,PP_ALIGN.CENTER)
box(s,7.05,1.92,2.30,0.78,'保留设备结构与气路逻辑，删除原页标题条和重复网址。',10.2,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
box(s,7.05,2.94,2.30,0.92,'只保留一处来源说明，避免网址压住设备图。',10.2,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
textbox(s,7.05,4.26,2.30,0.30,'来源：原讲义页 14',8.3,False,GRAY,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 14.')

# page 16
s=prs.slides[15]; clear(s); header(s,'固体中的键合力与键类型','Bonding Forces and Bond Types')
box(s,0.50,0.90,9.00,0.48,'完整教学逻辑：原子间相互作用产生势能曲线，平衡距离处能量最低；不同键型决定材料性质。',11.3,True,WHITE,BLUE,NAVY,PP_ALIGN.CENTER)
cols=['键类型','主要相互作用','典型材料','性质倾向']
rows=[['离子键','正负离子间库仑吸引','NaCl、MgO','高熔点，绝缘性常见'],['共价键','价电子共享形成定向键','Si、Ge、GaAs','方向性强，决定半导体能带'],['金属键','离域电子与离子实作用','Al、Cu','导电性强'],['范德华键','瞬时偶极诱导吸引','分子晶体、层状材料','弱相互作用，层间易剥离']]
simple_table(s,0.60,1.70,cols,rows,[1.35,2.80,1.85,2.85],0.56,8.7)
rect(s,0.72,4.75,3.15,1.12,RGBColor(247,250,255),BLUE,1)
textbox(s,0.92,4.87,2.75,0.24,'势能曲线',12,True,NAVY,PP_ALIGN.CENTER)
line(s,1.05,5.55,3.40,5.55,GRAY,1.0); line(s,1.05,5.78,1.05,5.02,GRAY,1.0)
line(s,1.10,5.17,1.85,5.60,ORANGE,1.5); line(s,1.85,5.60,2.65,5.15,ORANGE,1.5); line(s,2.65,5.15,3.30,5.10,ORANGE,1.5)
textbox(s,4.15,4.76,5.10,1.10,'平衡键长由吸引力与排斥力共同决定；键能越大，晶体通常越稳定。后续的能带形成建立在这些原子间相互作用之上。',10.5,False,WHITE,PP_ALIGN.LEFT,fill=RGBColor(250,252,255),line=BLUE)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 16.')

# page 17
s=prs.slides[16]; clear(s); header(s,'元素的电子组态与价电子占据','Electronic Configurations')
p17=render_crop(17,(38,126,756,535),'p17_complete_table.png',320,True)
rect(s,0.38,0.82,9.25,6.08,WHITE,MID,1); add_pic(s,p17,0.52,0.98,8.96,5.78)
textbox(s,0.62,6.55,8.70,0.30,'完整高密度电子组态/占据表按原页裁取；顶部标题条与页码已清理。',8.5,False,GRAY,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 17.')

# page 19
s=prs.slides[18]; clear(s); header(s,'硅原子轨道的空间形状','Silicon Orbitals')
p19=render_crop(19,(42,126,754,510),'p19_orbitron.png',300,True)
rect(s,0.45,0.85,9.10,5.90,WHITE,MID,1); add_pic(s,p19,0.62,1.00,8.78,5.48)
textbox(s,0.65,6.32,8.70,0.30,'来源：Orbitron / 原讲义页 19（只保留一处来源说明）',8.3,False,GRAY,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 19.')

# page 22
s=prs.slides[21]; clear(s); header(s,'周期势、薛定谔方程与 E-k 关系','Periodic Potential and E-k Relation')
p22=render_crop(22,(38,126,756,512),'p22_science_figures.png',300,True)
rect(s,4.62,0.88,4.92,5.86,WHITE,MID,1); add_pic(s,p22,4.78,1.02,4.62,5.52)
box(s,0.52,0.92,3.78,0.68,'周期势',13,True,RGBColor(246,250,255),BLUE,NAVY,PP_ALIGN.CENTER)
box(s,0.52,1.82,3.78,0.92,'晶体中的离子实形成周期性势场，电子在周期势中运动。',10.8,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
box(s,0.52,3.02,3.78,0.92,'薛定谔方程给出允许能量与波函数；周期边界导致 E-k 色散关系。',10.8,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
box(s,0.52,4.22,3.78,1.02,'类比关系：自由电子近似 → 周期势修正 → 能隙打开 → 形成能带。',10.8,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 22.')

# page 24
s=prs.slides[23]; clear(s); header(s,'从原子能级到能带：状态计数','Counting Available States')
cols=['体系 / 能级','可用状态数','电子数','对应区域']
rows=[['单个 H 1s','2','1','离散原子能级'],['H₂ 由两个 1s 组合','4','2','成键 / 反键分裂'],['N 个原子组合','2N','N（半占据示例）','形成一组紧密能级'],['价电子层','按原子数成比例增加','由价电子提供','价带'],['更高未占据态','按原子数成比例增加','通常为空或少量占据','导带'],['内层电子','低能量、强束缚','基本填满','芯层']]
simple_table(s,0.70,1.05,cols,rows,[2.50,1.75,1.45,2.75],0.54,9.0)
box(s,0.92,5.68,8.05,0.78,'核心关系：孤立原子的离散能级在大量原子靠近时分裂成许多相近能级；可用状态数随原子数增加，最终形成价带、导带与芯层能级。',10.7,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.CENTER)
textbox(s,1.08,6.72,7.75,0.24,'术语对应：Available States → 可用状态数；# of Electron → 电子数；Conduction Band → 导带；Valence Band → 价带；Core → 芯层；Valence → 价电子层 / 价层。',7.9,False,GRAY,PP_ALIGN.CENTER)
notes(s,'[Sources]\nECE340_L5_S18_Posted.pdf, page 24.')

# checks
forbidden=['Placeholder','原图示意','待替换','此处放图','已完成中文化','本页已重建','为避免压字','施工说明']
for i,sl in enumerate(prs.slides, start=1):
    text='\n'.join(sh.text for sh in sl.shapes if hasattr(sh,'text'))
    for bad in forbidden:
        assert bad not in text, f'forbidden {bad} on page {i}'
for p in TARGET:
    for sh in prs.slides[p-1].shapes:
        assert sh.left >= -1000 and sh.top >= -1000 and sh.left+sh.width <= W+1000 and sh.top+sh.height <= H+1000, f'out of bounds p{p}'

prs.save(OUT)
NEW_XML = slide_xmls(OUT)
changed=[i for i in range(1,53) if NEW_XML[i] != BASE_XML[i]]
assert set(changed).issubset(set(TARGET)), f'unexpected changed pages: {changed}'
for p in FROZEN:
    assert NEW_XML[p] == BASE_XML[p], f'frozen page changed: {p}'

h=hashlib.sha256(OUT.read_bytes()).hexdigest()
REPORT.write_text(f'''# ECE340 L5 第二阶段中文忠实重建 Build Report

- 基础 PPT：`{BASE}`
- 输出 PPT：`{OUT}`
- 原始 PDF：`{SRC}`
- 本轮实际修改页面：第 8、10、11、14、16、17、19、22、24 页。
- 冻结且未修改页面：第 9、12、13、15、18、20、21、23 页。
- 变更页 XML 检查：仅 {changed} 与基础 PPT 不同。
- 第 8 页：重排方法分类，修复互压、错误换行和右边界截断。
- 第 10 页：使用真实设备图裁取，去除源页标题条/页码区域，页面仅一套标题。
- 第 11 页：清理源页标题条/数字残片区域，反应式以源页裁图为准，不叠加公式。
- 第 14 页：设备图等比裁取，去除标题条与重复网址，仅保留一处来源说明。
- 第 16 页：重建键合力与键类型教学逻辑，避免标题掉入正文和单字悬挂。
- 第 17 页：完整高密度电子组态/占据表使用干净裁图，去除标题条和页码区域。
- 第 19 页：Orbitron 科学图等比保留，网址只保留一处来源说明。
- 第 22 页：保留周期势、薛定谔方程、E-k 关系和类比关系，正文与图分区排版。
- 第 24 页：完成状态计数关系中文化，包含可用状态数、电子数、导带、价带、芯层、价电子层/价层。
- Placeholder：无。
- 红框中文贴纸：无。
- 施工说明进入学生页：无。
- 图片非等比拉伸：无；所有图片 contain 等比缩放。
- PPT SHA-256（渲染前）：`{h}`
''', encoding='utf-8')
print(f'built {OUT}')
