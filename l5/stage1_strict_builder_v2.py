from pathlib import Path
import zipfile, hashlib
import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

ROOT = Path("l5")
BASE = ROOT / "ECE340_L5_S18_Posted_中文忠实重建_视觉返修版_第8-24页.pptx"
SRC = ROOT / "stage1_source_reference/ECE340_L5_S18_Posted.pdf"
OUT = ROOT / "ECE340_L5_S18_Posted_中文忠实重建_第一阶段严格返修_第9_12_18_20_23页.pptx"
ASSET = ROOT / "stage1_strict_assets"
REPORT = ROOT / "BUILD_REPORT_08_24_STAGE1_STRICT.md"
TARGET = [9, 12, 18, 20, 23]
ASSET.mkdir(exist_ok=True)

prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52 and doc.page_count == 52
W, H = prs.slide_width, prs.slide_height
assert abs(W/914400 - 10.0) < 0.01 and abs(H/914400 - 7.5) < 0.01

NAVY=RGBColor(30,50,82); BLUE=RGBColor(38,82,132); TEAL=RGBColor(45,119,122)
GOLD=RGBColor(208,153,55); GRAY=RGBColor(92,98,108); BLACK=RGBColor(28,28,30)
WHITE=RGBColor(255,255,255); PALE=RGBColor(240,246,252); PALETEAL=RGBColor(237,248,247)
LIGHT=RGBColor(250,252,255); MID=RGBColor(204,216,230)

def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)

def textbox(slide,x,y,w,h,text,size=12,bold=False,color=BLACK,align=PP_ALIGN.LEFT,
            fill=None,line=None,margin=3):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=line; sh.line.width=Pt(0.9)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(margin)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name="Noto Sans CJK SC"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def box(slide,x,y,w,h,text="",size=11,bold=False,fill=LIGHT,line=BLUE,color=BLACK,align=PP_ALIGN.LEFT):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(1.0)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Pt(6); tf.margin_top=tf.margin_bottom=Pt(3)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name="Noto Sans CJK SC"; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def rect(slide,x,y,w,h,fill=WHITE,line=MID,width=1):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh

def oval(slide,x,y,w,h,fill=PALE,line=BLUE,width=0.8):
    sh=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh

def connector(slide,x1,y1,x2,y2,color=GRAY,width=1):
    sh=slide.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    sh.line.color.rgb=color; sh.line.width=Pt(width); return sh

def right_arrow(slide,x,y,w,h,fill=TEAL):
    sh=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); return sh

def header(slide,zh,en):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    textbox(slide,0.38,0.06,6.15,0.42,zh,18,True,WHITE)
    textbox(slide,6.45,0.10,3.15,0.32,en,8.5,False,RGBColor(224,233,243),PP_ALIGN.RIGHT)

def render_crop(page_no,coords,name,dpi=240):
    pix=doc[page_no-1].get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),clip=fitz.Rect(*coords),alpha=False)
    out=ASSET/name; pix.save(out); return out

def add_picture_contain(slide,path,x,y,w,h):
    with Image.open(path) as im: iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    return slide.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

# p9
s=prs.slides[8]; clear(s); header(s,"硅的气相外延（VPE）","Silicon Vapor Phase Epitaxy")
box(s,0.42,0.75,4.35,0.48,"SiCl₄ + 2H₂ ⇌ Si + 4HCl",14,False,WHITE,BLUE,BLACK,PP_ALIGN.CENTER)
box(s,4.98,0.75,3.72,0.48,"SiH₄ → Si + 2H₂",14,False,WHITE,BLUE,BLACK,PP_ALIGN.CENTER)
photo=render_crop(9,(36,126,225,288),"p09_equipment.png")
rect(s,0.42,1.40,2.55,4.72,WHITE,MID,1)
add_picture_contain(s,photo,0.54,1.52,2.31,4.08)
textbox(s,0.62,5.68,2.15,0.28,"VPE 外延设备",9.5,True,NAVY,PP_ALIGN.CENTER)
rect(s,3.15,1.40,6.43,4.72,RGBColor(248,251,255),BLUE,1.1)
textbox(s,5.25,1.50,2.05,0.28,"石英反应腔",10.5,True,NAVY,PP_ALIGN.CENTER)
textbox(s,3.30,1.83,0.88,0.25,"气体入口",8.5,True,TEAL,PP_ALIGN.CENTER)
right_arrow(s,3.36,2.12,0.85,0.21,TEAL)
for xx in (4.25,4.95,5.65,6.35,7.05):
    right_arrow(s,xx,2.14,0.55,0.16,TEAL)
rect(s,4.24,1.90,0.11,1.08,RGBColor(236,224,191),GOLD,0.8)
textbox(s,3.86,2.97,0.90,0.24,"挡流板",8,True,GRAY,PP_ALIGN.CENTER)
connector(s,4.25,2.98,4.29,2.55,GRAY,0.8)
rect(s,5.25,3.72,2.92,0.20,RGBColor(112,121,134),GRAY,0.8)
for xx in (5.38,6.05,6.72,7.39):
    oval(s,xx,3.38,0.48,0.18,RGBColor(193,224,246),BLUE,0.7)
textbox(s,6.02,3.02,1.28,0.25,"硅片 / 晶圆",8.5,True,NAVY,PP_ALIGN.CENTER)
connector(s,6.65,3.25,6.95,3.43,NAVY,0.8)
textbox(s,6.02,3.98,1.40,0.24,"基座（Susceptor）",8,True,GRAY,PP_ALIGN.CENTER)
rect(s,6.42,3.92,0.47,0.86,RGBColor(204,210,216),GRAY,0.8)
textbox(s,7.10,4.26,1.30,0.24,"支座 / Pedestal",7.8,True,GRAY)
connector(s,7.12,4.38,6.88,4.35,GRAY,0.8)
for i in range(4):
    right_arrow(s,5.18+i*0.55,5.02,0.42,0.15,GOLD)
textbox(s,5.52,5.28,1.24,0.25,"RF 加热",8.5,True,GOLD,PP_ALIGN.CENTER)
right_arrow(s,8.18,2.12,0.95,0.21,TEAL)
textbox(s,8.32,1.82,0.80,0.25,"排气口",8.5,True,TEAL,PP_ALIGN.CENTER)
textbox(s,3.55,5.66,5.55,0.30,"气流沿反应腔流过加热晶圆表面，再由排气口排出。",9.5,False,BLACK,PP_ALIGN.CENTER)
notes(s,"VPE 反应气体沿反应腔通过加热晶圆表面。两条反应式严格按原页保留。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 9.")

# p12
s=prs.slides[11]; clear(s); header(s,"分子束外延（MBE）","Molecular Beam Epitaxy")
rect(s,0.40,0.88,4.50,5.95,RGBColor(248,251,255),BLUE,1.1)
textbox(s,1.48,1.00,2.35,0.28,"超高真空生长腔",11.5,True,NAVY,PP_ALIGN.CENTER)
rect(s,1.88,1.70,1.58,0.24,RGBColor(183,216,243),BLUE,0.8)
textbox(s,1.86,1.98,1.62,0.25,"GaAs 衬底",8.8,True,NAVY,PP_ALIGN.CENTER)
rect(s,2.28,1.42,0.72,0.14,RGBColor(188,192,198),GRAY,0.7)
textbox(s,3.12,1.37,0.92,0.24,"衬底支架",7.5,False,GRAY)
connector(s,3.14,1.49,2.98,1.52,GRAY,0.7)
sources=[("Si",0.68),("Al",1.34),("Ga",2.00),("As",2.66),("Be",3.32)]
for lab,x in sources:
    rect(s,x,5.62,0.40,0.46,RGBColor(239,241,244),GRAY,0.8)
    textbox(s,x,5.67,0.40,0.20,lab,7.5,True,NAVY,PP_ALIGN.CENTER)
    rect(s,x-0.01,5.23,0.42,0.06,RGBColor(150,156,164),GRAY,0.2)
    connector(s,x+0.20,5.20,2.67,2.22,GOLD,1.0)
textbox(s,0.66,6.16,3.30,0.25,"束源：Si / Al / Ga / As / Be",8,True,GRAY,PP_ALIGN.CENTER)
textbox(s,3.50,4.92,0.88,0.24,"快门",8,True,GRAY,PP_ALIGN.CENTER)
connector(s,3.52,5.08,3.33,5.25,GRAY,0.7)
textbox(s,0.90,3.34,1.10,0.24,"分子束",8.5,True,GOLD,PP_ALIGN.CENTER)
connector(s,1.82,3.45,2.25,3.12,GOLD,0.8)
micro=render_crop(12,(390,126,720,535),"p12_right_source.png")
rect(s,5.10,0.88,4.50,5.95,WHITE,MID,1)
add_picture_contain(s,micro,5.24,1.02,4.22,5.65)
notes(s,"左图保留 Si / Al / Ga / As / Be 束源、快门、分子束与 GaAs 衬底的对应关系。右图直接使用原页真实图像裁取，不添加中文贴纸；图内 4×4 / GaAs substrate / 10 nm / <100> 按原图保留。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 12.")

# p18
s=prs.slides[17]; clear(s); header(s,"硅的芯层电子与价电子","Silicon Core and Valence Electrons")
p18=render_crop(18,(54,126,738,535),"p18_complete_source_body.png",260)
rect(s,0.40,0.82,9.20,6.18,WHITE,MID,1)
add_picture_contain(s,p18,0.55,0.97,8.90,5.88)
notes(s,"本页两幅高信息密度原图完整保留，不在英文原图上叠加中文框。讲解重点：+14；1s / 2s / 2p / 3s / 3p；芯层电子、价电子、价轨道、首次激发轨道、价能级与电离/零能级的关系。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 18.")

# p20
s=prs.slides[19]; clear(s); header(s,"sp³ 杂化与四面体几何","sp³ Hybridization")
p20=render_crop(20,(42,126,754,535),"p20_formula_geometry_source.png",270)
rect(s,0.40,0.82,9.20,6.18,WHITE,MID,1)
add_picture_contain(s,p20,0.55,0.97,8.90,5.88)
notes(s,"原页四条 sp³ 杂化公式与 109.5° 四面体几何直接使用源页裁图保留，不重新誊写、不改系数。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 20.")

# p23
s=prs.slides[22]; clear(s); header(s,"氢分子的成键与反键轨道","Bonding / Antibonding Orbitals in H₂")
box(s,0.40,0.86,2.38,2.12,
    "元素氢：1s¹\nH #1: 2 states, 1 electron\nH #2: 2 states, 1 electron\nH₂: 4 states, 2 electrons",
    9.5,False,PALE,BLUE,BLACK)
box(s,0.40,3.24,2.38,1.20,
    "两个 H 1s 原子轨道组合后形成：\n低能成键轨道 + 高能反键轨道。",
    9.2,True,WHITE,BLUE,BLACK,PP_ALIGN.CENTER)
box(s,0.40,4.70,2.38,1.52,
    "H₂ 的 2 个电子在基态优先占据低能成键轨道，形成稳定 H–H 键。",
    9.2,False,PALETEAL,TEAL,BLACK,PP_ALIGN.CENTER)
rect(s,2.98,0.86,4.05,5.96,RGBColor(251,253,255),MID,1)
textbox(s,3.95,0.99,2.08,0.26,"分子轨道能级关系",10.5,True,NAVY,PP_ALIGN.CENTER)
connector(s,3.30,3.78,4.08,3.78,GRAY,1.5)
connector(s,5.90,3.78,6.68,3.78,GRAY,1.5)
textbox(s,3.25,3.40,0.88,0.24,"H #1 1s",7.7,True,GRAY,PP_ALIGN.CENTER)
textbox(s,5.86,3.40,0.88,0.24,"H #2 1s",7.7,True,GRAY,PP_ALIGN.CENTER)
connector(s,4.42,2.10,5.56,2.10,NAVY,2)
textbox(s,4.24,1.69,1.50,0.26,"反键 σ*1s",8.5,True,NAVY,PP_ALIGN.CENTER)
connector(s,4.42,5.16,5.56,5.16,TEAL,2)
textbox(s,4.24,5.30,1.50,0.26,"成键 σ1s",8.5,True,TEAL,PP_ALIGN.CENTER)
for a,b,c,d in [(4.08,3.78,4.42,2.10),(4.08,3.78,4.42,5.16),(5.90,3.78,5.56,2.10),(5.90,3.78,5.56,5.16)]:
    connector(s,a,b,c,d,GRAY,0.9)
textbox(s,4.75,4.70,0.30,0.30,"↑",15,True,TEAL,PP_ALIGN.CENTER)
textbox(s,5.06,4.70,0.30,0.30,"↓",15,True,TEAL,PP_ALIGN.CENTER)
up=s.shapes.add_shape(MSO_SHAPE.UP_ARROW,Inches(6.50),Inches(1.63),Inches(0.20),Inches(3.88))
up.fill.solid(); up.fill.fore_color.rgb=NAVY; up.line.fill.background()
textbox(s,6.05,1.18,0.68,0.38,"Higher\nEnergy",6.7,True,NAVY,PP_ALIGN.CENTER)
textbox(s,6.05,5.60,0.68,0.38,"Lower\nEnergy",6.7,True,TEAL,PP_ALIGN.CENTER)
rect(s,7.22,0.86,2.38,5.96,WHITE,MID,1)
textbox(s,7.40,1.00,2.02,0.26,"轨道与电子密度",9.5,True,NAVY,PP_ALIGN.CENTER)
oval(s,7.55,1.72,0.70,0.70,RGBColor(229,239,250),BLUE,0.8)
oval(s,8.60,1.72,0.70,0.70,RGBColor(247,236,216),GOLD,0.8)
textbox(s,7.78,1.93,0.24,0.22,"+",9,True,NAVY,PP_ALIGN.CENTER)
textbox(s,8.83,1.93,0.24,0.22,"−",9,True,GOLD,PP_ALIGN.CENTER)
rect(s,8.39,1.64,0.035,0.90,GRAY,GRAY,0)
textbox(s,7.38,2.56,2.05,0.60,"反键：两核间有节点\n电子不位于两核之间",7.5,True,NAVY,PP_ALIGN.CENTER)
oval(s,7.57,4.08,0.86,0.73,RGBColor(220,240,238),TEAL,0.8)
oval(s,8.43,4.08,0.86,0.73,RGBColor(220,240,238),TEAL,0.8)
oval(s,8.10,4.17,0.68,0.55,RGBColor(195,227,224),TEAL,0.6)
textbox(s,7.38,4.98,2.05,0.60,"成键：两核间电子密度增加\n系统能量降低",7.5,True,TEAL,PP_ALIGN.CENTER)
notes(s,"H #1 与 H #2 的 1s 原子轨道组合为低能成键轨道和高能反键轨道。成键轨道在两核之间电子密度增加；反键轨道在两核之间形成节点。H₂ 的两个电子在基态占据成键轨道。\n[Sources]\nECE340_L5_S18_Posted.pdf, page 23.")

for p in TARGET:
    for sh in prs.slides[p-1].shapes:
        assert sh.left >= 0 and sh.top >= 0
        assert sh.left+sh.width <= W+1000,(p,sh.name,"x overflow", (sh.left+sh.width)/914400, W/914400)
        assert sh.top+sh.height <= H+1000,(p,sh.name,"y overflow")

prs.save(OUT)

with zipfile.ZipFile(BASE) as zb, zipfile.ZipFile(OUT) as zo:
    for p in range(1,53):
        a=zb.read(f"ppt/slides/slide{p}.xml"); b=zo.read(f"ppt/slides/slide{p}.xml")
        if p in TARGET: assert a!=b,f"target {p} unchanged"
        else: assert a==b,f"non-target page {p} changed"
    for p in range(1,53):
        n=f"ppt/notesSlides/notesSlide{p}.xml"
        if p not in TARGET and n in zb.namelist() and n in zo.namelist():
            assert zb.read(n)==zo.read(n),f"non-target notes {p} changed"

check=Presentation(OUT)
banned=["原设备照片","右侧显微图","原图示意","待替换图片","此处保留原图","本页已重建","已删除残留",
        "保留原始信息","为避免 PowerPoint 压字","本页采用中文示意","右图保留","已完成中文化","已重做","已改为"]
for p in TARGET:
    texts=[]
    for sh in check.slides[p-1].shapes:
        assert sh.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER,(p,sh.name)
        if getattr(sh,"has_text_frame",False): texts.append(sh.text)
    joined="\n".join(texts)
    for term in banned: assert term not in joined,(p,term)

t9="\n".join(sh.text for sh in check.slides[8].shapes if getattr(sh,"has_text_frame",False))
for term in ["SiCl₄ + 2H₂ ⇌ Si + 4HCl","SiH₄ → Si + 2H₂","石英反应腔","气体入口","硅片 / 晶圆","基座（Susceptor）","支座 / Pedestal","RF 加热","排气口"]:
    assert term in t9,term
t12="\n".join(sh.text for sh in check.slides[11].shapes if getattr(sh,"has_text_frame",False))
for term in ["Si / Al / Ga / As / Be","GaAs 衬底","快门","分子束"]: assert term in t12,term
t23="\n".join(sh.text for sh in check.slides[22].shapes if getattr(sh,"has_text_frame",False))
for term in ["H #1: 2 states, 1 electron","H₂: 4 states, 2 electrons","Higher","Lower","成键","反键"]: assert term in t23,term

src20=doc[19].get_text("text").replace("\n"," ")
for seq in ["ψ 1 = 1","ψ s +ψ px +ψ py +ψ pz","ψ 2 = 1","ψ s −ψ px −ψ py +ψ pz",
            "ψ 3 = 1","ψ s +ψ px −ψ py −ψ pz","ψ 4 = 1","ψ s −ψ px +ψ py −ψ pz","109.5° Bond"]:
    assert seq in src20,seq

digest=hashlib.sha256(OUT.read_bytes()).hexdigest()
REPORT.write_text(
f"""# ECE340 L5 第一阶段严格视觉返修 Build Report

- 基准 PPT：`{BASE.as_posix()}`
- 原始讲义 PDF：`{SRC.as_posix()}`
- 阶段输出 PPT：`{OUT.as_posix()}`
- 页面尺寸：10.0 × 7.5 inch（4:3）；所有目标页 shape 均通过版心边界断言。
- 实际修改页：9、12、18、20、23。
- 第 8–24 页中保持未改：8、10、11、13、14、15、16、17、19、21、22、24。
- 其他保持未改：1–7、25–52。
- 页 9：保留原页真实设备照片；重建 VPE 气流、反应腔、晶圆、基座、支座、排气口、RF 加热关系；两条反应式逐字断言。
- 页 12：重建 Si/Al/Ga/As/Be 束源—快门—分子束—衬底关系；右侧使用原页真实裁图。
- 页 18：两幅高密度关系图使用原页正文干净裁图完整保留；中文讲解置备注。
- 页 20：四条 sp³ 公式与 109.5° 直接使用原页正文裁图，不重新誊写或换公式。
- 页 23：重建状态数、原子 1s—分子轨道能级、成键/反键、Higher/Lower Energy 与电子密度对应。
- Placeholder：无；目标页已检查 shape_type。
- 红框中文贴纸：无。
- 施工说明：无；已执行禁词扫描。
- 科学内容/公式：未修改；页 20 原始公式文本已从源 PDF 逐项校验，页面直接使用源裁图。
- 非目标页 slide XML：与基准逐页 byte-for-byte 相同。
- PPT SHA-256（渲染前）：`{digest}`
""",encoding="utf-8")
print("built",OUT)
