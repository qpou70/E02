from pathlib import Path
import zipfile, hashlib
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段返修_第8_10_11_14_16_17_19_22_24页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND1_第8_10_11_14_16_17_19_22_24页.pptx'
ASSET = ROOT / 'stage2_round1_assets'
REPORT = ROOT / 'BUILD_REPORT_STAGE2_ROUND1.md'
TARGET = [8,10,11,14,16,17,19,22,24]
FROZEN = [9,12,13,15,18,20,21,23]
FORBIDDEN = ['源页','原页裁取','去除标题条','已清理','页面只保留','只保留一处','等比裁取','不叠加公式','真实设备照片','已完成中文化','本页已重建','保留原图','施工','待替换','占位','Placeholder','此处放图']
ASSET.mkdir(exist_ok=True)

prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52 and doc.page_count == 52
W, H = prs.slide_width, prs.slide_height

NAVY=RGBColor(30,50,82); BLUE=RGBColor(38,82,132); TEAL=RGBColor(45,119,122)
GOLD=RGBColor(208,153,55); GRAY=RGBColor(92,98,108); BLACK=RGBColor(28,28,30)
WHITE=RGBColor(255,255,255); PALE=RGBColor(240,246,252); MID=RGBColor(204,216,230)
LIGHT=RGBColor(250,252,255); GREEN=RGBColor(73,130,84); ORANGE=RGBColor(209,122,42)
RED=RGBColor(160,58,58)

def slide_xmls(pptx):
    out={}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1,53):
            out[i]=z.read(f'ppt/slides/slide{i}.xml')
    return out

BASE_XML = slide_xmls(BASE)

def clear(slide):
    for sh in list(slide.shapes):
        sh.element.getparent().remove(sh.element)

def textbox(slide,x,y,w,h,text,size=12,bold=False,color=BLACK,align=PP_ALIGN.LEFT,fill=None,line=None,margin=3):
    sh=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb=fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb=line; sh.line.width=Pt(0.8)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Pt(margin)
    tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    for k, part in enumerate(str(text).split('\n')):
        if k:
            p=tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def rect(slide,x,y,w,h,fill=WHITE,line=MID,width=1):
    sh=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(width); return sh

def box(slide,x,y,w,h,text,size=11,bold=False,fill=LIGHT,line=BLUE,color=BLACK,align=PP_ALIGN.LEFT,margin=5):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(0.9)
    tf=sh.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Pt(margin); tf.margin_top=tf.margin_bottom=Pt(4); tf.vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.alignment=align
    for k, part in enumerate(str(text).split('\n')):
        if k:
            p=tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=part; r.font.name='Noto Sans CJK SC'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color
    return sh

def line(slide,x1,y1,x2,y2,color=GRAY,width=1):
    sh=slide.shapes.add_connector(1,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    sh.line.color.rgb=color; sh.line.width=Pt(width); return sh

def arrow(slide,x,y,w,h,fill=TEAL):
    sh=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); return sh

def header(slide,zh,en):
    bar=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.58))
    bar.fill.solid(); bar.fill.fore_color.rgb=NAVY; bar.line.fill.background()
    textbox(slide,0.38,0.06,6.35,0.42,zh,18,True,WHITE)
    textbox(slide,6.48,0.10,3.10,0.32,en,8.5,False,RGBColor(224,233,243),PP_ALIGN.RIGHT)

def notes(slide,text):
    slide.notes_slide.notes_text_frame.text=text

def render_crop(page_no,coords,name,dpi=275,mask_right_page_number=True):
    pix=doc[page_no-1].get_pixmap(matrix=fitz.Matrix(dpi/72,dpi/72),clip=fitz.Rect(*coords),alpha=False)
    out=ASSET/name; pix.save(out)
    if mask_right_page_number:
        im=Image.open(out).convert('RGB'); d=ImageDraw.Draw(im)
        # remove possible page number at lower-right without touching the figure body
        d.rectangle([im.width-90, im.height-42, im.width, im.height], fill='white')
        im.save(out)
    return out

def add_pic(slide,path,x,y,w,h):
    with Image.open(path) as im:
        iw,ih=im.size
    k=min(w/iw,h/ih); pw,ph=iw*k,ih*k
    return slide.shapes.add_picture(str(path),Inches(x+(w-pw)/2),Inches(y+(h-ph)/2),Inches(pw),Inches(ph))

def simple_table(slide,x,y,cols,rows,widths,row_h=0.42,font=8.6):
    cx=x
    for j,c in enumerate(cols):
        box(slide,cx,y,widths[j],row_h,c,font,True,PALE,BLUE,NAVY,PP_ALIGN.CENTER,margin=2); cx+=widths[j]
    for i,row in enumerate(rows):
        cx=x
        for j,c in enumerate(row):
            fill=RGBColor(252,253,255) if i%2==0 else WHITE
            box(slide,cx,y+row_h*(i+1),widths[j],row_h,str(c),font,False,fill,MID,BLACK,PP_ALIGN.CENTER if j else PP_ALIGN.LEFT,margin=2); cx+=widths[j]

# Page 8: Epitaxial Growth Methods only
s=prs.slides[7]; clear(s); header(s,'外延生长方法','Epitaxial Growth Methods')
box(s,0.65,0.92,8.70,0.48,'外延生长是在晶体基底上生长具有受控取向和组成的薄层。',12.0,True,WHITE,BLUE,NAVY,PP_ALIGN.CENTER)
methods=[
    ('液相外延（LPE）','Liquid Phase Epitaxy'),
    ('气相外延（VPE）','Vapor Phase Epitaxy\n（三）氯化物与氢化物 VPE'),
    ('分子束外延（MBE）','Molecular Beam Epitaxy'),
    ('化学气相沉积（CVD）','Chemical Vapor Deposition'),
    ('金属有机化学气相沉积（MOCVD）','Metalorganic Chemical Vapor Deposition')]
y=1.65
for idx,(zh,en) in enumerate(methods, start=1):
    box(s,0.98,y,0.48,0.48,str(idx),14,True,PALE,BLUE,NAVY,PP_ALIGN.CENTER)
    box(s,1.62,y,3.05,0.48,zh,12.0,True,WHITE,BLUE,NAVY,PP_ALIGN.CENTER)
    box(s,4.90,y,3.95,0.48,en,9.5,False,WHITE,MID,BLACK,PP_ALIGN.CENTER)
    y += 0.78
box(s,1.08,5.85,7.78,0.58,'本页列出后续讨论的外延方法：LPE、VPE、MBE、CVD 与 MOCVD。',11.5,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.CENTER)
notes(s,'Page 8 rebuilt from ECE340_L5_S18_Posted.pdf, page 8: Epitaxial Growth Methods list only.')

# Page 10: MOCVD early systems
s=prs.slides[9]; clear(s); header(s,'MOCVD：早期系统','MOCVD: Early Systems')
p10=render_crop(10,(36,126,756,515),'round1_p10_mocvd_early.png',300,True)
rect(s,0.52,0.82,8.96,5.76,WHITE,MID,1); add_pic(s,p10,0.68,0.98,8.64,5.30)
textbox(s,0.90,6.50,8.20,0.32,'早期 MOCVD 系统由复杂的气路、流量控制和反应腔组成。',10.5,False,GRAY,PP_ALIGN.CENTER)
notes(s,'Page 10 rebuilt from ECE340_L5_S18_Posted.pdf, page 10: MOCVD Early Systems equipment photo.')

# Page 11: MOCVD today, photos and reaction
s=prs.slides[10]; clear(s); header(s,'MOCVD：现代系统','MOCVD: Today')
box(s,1.05,0.82,7.90,0.48,'(CH₃)₃Ga + AsH₃  →  GaAs + 3CH₄',15,False,WHITE,BLUE,BLACK,PP_ALIGN.CENTER)
p11=render_crop(11,(36,165,756,520),'round1_p11_mocvd_today_photos.png',300,True)
rect(s,0.52,1.46,8.96,5.38,WHITE,MID,1); add_pic(s,p11,0.66,1.62,8.68,5.05)
notes(s,'Page 11 rebuilt from ECE340_L5_S18_Posted.pdf, page 11: MOCVD Today photos and reaction equation.')

# Page 14: Production MBE Reactor
s=prs.slides[13]; clear(s); header(s,'生产型 MBE 反应器','Production MBE Reactor')
p14=render_crop(14,(36,126,756,505),'round1_p14_production_mbe.png',300,True)
rect(s,0.52,0.86,8.96,5.78,WHITE,MID,1); add_pic(s,p14,0.66,1.02,8.68,5.36)
textbox(s,0.80,6.58,8.40,0.24,'www.mbe-komponenten.de',8.2,False,GRAY,PP_ALIGN.CENTER)
notes(s,'Page 14 rebuilt from ECE340_L5_S18_Posted.pdf, page 14: Production MBE Reactor.')

# Page 16: Bond types, faithfully restoring four categories
s=prs.slides[15]; clear(s); header(s,'键合类型','Bond Types')
box(s,0.58,0.88,4.18,2.58,'离子键 / Ionic Bonding\n\n电子转移。\n库仑吸引与原子核排斥在平衡距离处达到平衡。',11.2,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT)
box(s,5.22,0.88,4.18,2.58,'共价键 / Covalent Bonding\n\n电子共享。\n典型材料：Si、Ge、C。',11.2,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT)
box(s,0.58,3.72,4.18,2.58,'混合离子-共价键 / Mixed Ionic-Covalent Bonding\n\n由电负性差异产生极性共价键。\n典型材料：GaAs、InP、GaN。',10.5,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT)
box(s,5.22,3.72,4.18,2.58,'金属键 / Metallic Bonding\n\n正离子实处于电子海中。\n典型于价电子数不超过 3 的原子。',11.0,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT)
notes(s,'Page 16 rebuilt from ECE340_L5_S18_Posted.pdf, page 16: Bond Types.')

# Page 17: only remove bottom maker note; keep table and original relationships
s=prs.slides[16]
for sh in list(s.shapes):
    txt = getattr(sh, 'text', '') or ''
    if '完整高密度电子组态' in txt or any(bad in txt for bad in FORBIDDEN):
        sh.element.getparent().remove(sh.element)
notes(s,'Page 17: bottom non-teaching note removed; table, red arrow and highlighted relationships kept from the previous Stage 2 slide.')

# Page 19: generic orbital wave function shape, complete Orbitron figure
s=prs.slides[18]; clear(s); header(s,'原子轨道波函数的空间形状','Orbital Wave Function Shape')
p19=render_crop(19,(36,116,756,532),'round1_p19_orbitron_full.png',300,True)
rect(s,0.48,0.82,9.04,6.16,WHITE,MID,1); add_pic(s,p19,0.60,0.96,8.80,5.84)
notes(s,'Page 19 rebuilt from ECE340_L5_S18_Posted.pdf, page 19: Orbital Wave Function Shape.')

# Page 22: Chinese body text plus scientific figures
s=prs.slides[21]; clear(s); header(s,'周期势与 E-k 关系','Periodic Potential and E-k Relation')
# Use figure-only crop from the right side of original page; teaching body is rebuilt in Chinese on the left.
p22=render_crop(22,(352,128,756,518),'round1_p22_figures.png',300,True)
rect(s,5.42,0.88,4.04,5.98,WHITE,MID,1); add_pic(s,p22,5.54,1.02,3.80,5.66)
items=[
    ('定性理解','Streetman 教材和本课程给出的是能带形成的定性图像，用来说明为什么固体中会出现允许能带与禁带。'),
    ('理论范围','能带形成的完整量子理论超出本课程范围，本页只保留必要的物理背景和类比。'),
    ('周期势','在晶格周期势中求解薛定谔方程，电子波函数受到周期结构调制，从而产生能带。'),
    ('E-k 关系','求解结果给出电子能量 E 与晶体中动量波矢 k 的关系，即 E-k 关系。'),
    ('光学类比','可类比光学干涉滤光片或蝴蝶翅膀的衍射结构；差别在于尺度对应电子波长。')]
y=0.92
for title,body in items:
    box(s,0.58,y,4.55,0.88,f'{title}\n{body}',8.8,False,WHITE,BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
    y += 1.05
textbox(s,5.60,6.62,3.64,0.24,'E-k 图、晶格周期势与光学衍射类比',8.4,False,GRAY,PP_ALIGN.CENTER)
notes(s,'Page 22 rebuilt from ECE340_L5_S18_Posted.pdf, page 22: periodic potential, Schrodinger equation background, E-k relation and optical analogy.')

# Page 24: Si discrete levels to bands and state counting
s=prs.slides[23]; clear(s); header(s,'硅：从离散能级到能带的状态计数','Si States from Atomic Levels to Bands')
box(s,0.55,0.82,4.35,1.12,'原子靠近形成固体\n原子间吸引力与排斥力在合适的原子间距处达到平衡；离散原子能级随原子靠近而分裂，最终形成能带。',9.8,False,WHITE,BLUE,BLACK,PP_ALIGN.LEFT)
box(s,5.08,0.82,4.38,1.12,'能带结构\n分裂后的能级形成价带和导带，中间由禁带分隔；0 K 时价带充满，导带为空。',9.8,False,WHITE,BLUE,BLACK,PP_ALIGN.LEFT)
# simple diagram
line(s,1.10,2.38,3.52,2.38,GRAY,1.1); line(s,1.10,2.58,3.52,2.58,GRAY,1.1); line(s,1.10,2.78,3.52,2.78,GRAY,1.1)
textbox(s,1.05,2.03,2.55,0.22,'离散原子能级',9,True,NAVY,PP_ALIGN.CENTER)
arrow(s,3.75,2.38,0.62,0.28,TEAL)
rect(s,4.68,2.05,1.40,0.56,RGBColor(230,242,255),BLUE,1); textbox(s,4.78,2.16,1.20,0.24,'导带',11,True,NAVY,PP_ALIGN.CENTER)
rect(s,4.68,3.16,1.40,0.56,RGBColor(255,244,225),GOLD,1); textbox(s,4.78,3.27,1.20,0.24,'价带',11,True,NAVY,PP_ALIGN.CENTER)
textbox(s,4.72,2.72,1.30,0.24,'禁带',10,True,RED,PP_ALIGN.CENTER)
box(s,6.55,2.02,2.70,1.64,'Si 电子组态\n1s² 2s² 2p⁶ 3s² 3p²\n芯层：1s² 2s² 2p⁶\n价电子：3s² 3p²\n外层还缺 4 个电子才能填满。',9.0,False,RGBColor(250,252,255),BLUE,BLACK,PP_ALIGN.LEFT,margin=4)
# n=3 table
textbox(s,0.72,4.22,3.70,0.28,'n = 3 的状态计数',11.5,True,NAVY,PP_ALIGN.CENTER)
simple_table(s,0.72,4.56,['能级','可用状态数','电子数'],[['3s','2','2'],['3p','6','2'],['Total','8','4']],[1.05,1.45,1.05],0.38,8.6)
# N atoms table
textbox(s,5.05,4.22,4.15,0.28,'N 个 Si 原子形成晶体时',11.5,True,NAVY,PP_ALIGN.CENTER)
simple_table(s,4.70,4.56,['区域','可用状态数','电子数'],[['Total','N×8','N×4'],['导带','N×4','0（0 K）'],['价带','N×4','N×4（0 K）']],[1.25,1.55,1.35],0.38,8.6)
notes(s,'Page 24 rebuilt from ECE340_L5_S18_Posted.pdf, page 24: Si atomic configuration and state counting.')

# Remove any visible non-teaching maker words from target pages.
for p in TARGET:
    slide=prs.slides[p-1]
    for sh in list(slide.shapes):
        txt = getattr(sh,'text','') or ''
        if any(bad in txt for bad in FORBIDDEN):
            raise AssertionError(f'forbidden maker word on page {p}: {txt}')
        assert sh.left >= -1000 and sh.top >= -1000 and sh.left + sh.width <= W + 1000 and sh.top + sh.height <= H + 1000, f'out of bounds page {p}'

prs.save(OUT)
NEW_XML = slide_xmls(OUT)
changed=[i for i in range(1,53) if NEW_XML[i] != BASE_XML[i]]
assert set(changed).issubset(set(TARGET)), f'unexpected changed pages: {changed}'
for p in FROZEN:
    assert NEW_XML[p] == BASE_XML[p], f'frozen page changed: {p}'
for p in list(range(1,8)) + list(range(25,53)):
    assert NEW_XML[p] == BASE_XML[p], f'outside allowed range changed: {p}'

h=hashlib.sha256(OUT.read_bytes()).hexdigest()
REPORT.write_text(f'''# ECE340 L5 第二阶段视觉返修 ROUND1 Build Report

- 返修依据：`l5/SUPERVISOR_STAGE2_VISUAL_REVIEW_FEEDBACK_ROUND1.md`，反馈提交 `bc9f3b543667e3de5bb97da07b48526d43417127`。
- 当前返修基准 PPT：`{BASE}`
- 输出 PPT：`{OUT}`
- 原始 PDF：`{SRC}`
- 本轮实际修改页面：第 8、10、11、14、16、17、19、22、24 页。
- 冻结且确认未修改页面：第 9、12、13、15、18、20、21、23 页。
- 同时确认未修改页面：第 1–7、25–52 页。
- 变更页 XML 检查：仅 {changed} 与本轮基准 PPT 不同。

## 逐页修复摘要

- 第 8 页：恢复为 Epitaxial Growth Methods / 外延生长方法；删除上一版新增的 bulk growth、Czochralski、Float-zone、Bridgman；只保留 LPE、VPE、(Tri)Chloride & Hydride VPE、MBE、CVD、MOCVD。
- 第 10 页：恢复为 MOCVD: Early Systems / MOCVD：早期系统；保留早期 MOCVD 系统照片，删除全部学生页施工说明。
- 第 11 页：恢复为 MOCVD: Today / MOCVD：现代系统；恢复反应式 `(CH₃)₃Ga + AsH₃ → GaAs + 3CH₄`，保留现代系统相关照片，删除全部学生页施工说明。
- 第 14 页：恢复为 Production MBE Reactor / 生产型 MBE 反应器；保留生产型 MBE 设备图，避免 MBE/MOCVD 混淆。
- 第 16 页：恢复 Bond Types 四类：离子键、共价键、混合离子-共价键、金属键；删除上一版新增的范德华键与空白框。
- 第 17 页：仅删除底部制作说明，保留主电子组态/占据表、红箭头和高亮关系。
- 第 19 页：恢复通用 Orbital Wave Function Shape / 原子轨道波函数的空间形状；不再限定为硅原子；恢复 Orbitron 图底部完整性。
- 第 22 页：删除英文正文作为主要教学内容的做法；用中文完整重建 Streetman 定性理解、理论范围、周期势、薛定谔方程、E-k 关系和光学类比。
- 第 24 页：撤销 H/H₂ 通用框架；恢复 Si 电子组态、芯层/价电子、n=3 状态计数以及 N 个 Si 原子形成晶体后的导带/价带状态计数。

## 强制自检记录

- 学生 PPT 页面中无 `源页`、`原页裁取`、`去除标题条`、`已清理`、`页面只保留`、`只保留一处`、`等比裁取`、`不叠加公式`、`真实设备照片`、`已完成中文化`、`本页已重建`、`保留原图`、`施工`、`待替换`、`占位`、`Placeholder`、`此处放图`。
- 不使用红框中文贴纸。
- 图片按 contain 方式等比放置，不做非等比拉伸。
- PPT SHA-256（渲染前）：`{h}`
''', encoding='utf-8')
print(f'round1 stage2 repair built {OUT}')
