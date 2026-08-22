from pathlib import Path
import shutil
import zipfile
import hashlib
import subprocess
import re
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_最终候选版_第8-24页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_最终候选版_R1_第8-24页.pptx'
EV = ROOT / 'final_qa_r1_08_24'
RENDER = EV / 'rendered'
COMP = EV / 'comparison'
PDF_DIR = EV / 'final_pdf'
REPORT = ROOT / 'BUILD_REPORT_FINAL_QA_R1_08_24.md'
PAGES = list(range(8, 25))
VISUAL_ALLOWED = {13, 15, 20, 21}
NOTES_ALLOWED = {8, 9, 10, 11, 12, 14, 16, 17, 18, 19, 20, 22, 23, 24}
NOTES_FROZEN = {13, 15, 21}
for p in [EV, RENDER, COMP, PDF_DIR]:
    p.mkdir(parents=True, exist_ok=True)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def pptx_part_bytes(pptx: Path, part_prefix: str):
    data = {}
    with zipfile.ZipFile(pptx) as z:
        for name in z.namelist():
            if name.startswith(part_prefix):
                data[name] = z.read(name)
    return data


def slide_xmls(pptx: Path):
    data = {}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1, 53):
            data[i] = z.read(f'ppt/slides/slide{i}.xml')
    return data


def extract_notes(prs: Presentation, idx: int) -> str:
    try:
        return prs.slides[idx - 1].notes_slide.notes_text_frame.text or ''
    except Exception:
        return ''


def delete_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def iter_shapes(shapes):
    for sh in shapes:
        yield sh
        if hasattr(sh, 'shapes'):
            yield from iter_shapes(sh.shapes)


def text_of(shape):
    if hasattr(shape, 'text_frame') and shape.has_text_frame:
        return shape.text_frame.text or ''
    return ''


def set_shape_text(shape, text, size=24, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Noto Sans CJK SC'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)


def add_white_box(slide, x, y, w, h):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shp.line.color.rgb = RGBColor(255, 255, 255)
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = 'Noto Sans CJK SC'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)
    return box


def remove_bottom_page_number(slide, number: str, prs: Presentation):
    sw, sh = prs.slide_width, prs.slide_height
    for shape in list(slide.shapes):
        t = text_of(shape).strip()
        if t == number and shape.left > sw * 0.82 and shape.top > sh * 0.78:
            delete_shape(shape)
    # Backstop for page numbers that are embedded inside an image/text group.
    add_white_box(slide, 9.45, 7.00, 0.38, 0.30)


def replace_or_cover_caption_13(slide):
    done = False
    for shape in list(slide.shapes):
        t = text_of(shape)
        if 'Figure 1.17' in t or 'University of Texas at Austin Microelectronics Research Center' in t:
            set_shape_text(shape, '图 1.17　德克萨斯大学奥斯汀分校微电子研究中心的 MBE 设备', size=14, bold=False)
            done = True
    if not done:
        add_white_box(slide, 2.55, 5.88, 5.20, 0.52)
        add_text(slide, 2.55, 5.88, 5.20, 0.52, '图 1.17　德克萨斯大学奥斯汀分校微电子研究中心的 MBE 设备', size=14)


def replace_big_title(slide, old_substring: str, new_text: str, y_backstop: float):
    done = False
    for shape in list(slide.shapes):
        t = text_of(shape)
        if old_substring in t:
            set_shape_text(shape, new_text, size=40, bold=True)
            done = True
    if not done:
        add_white_box(slide, 1.00, y_backstop, 8.00, 1.00)
        add_text(slide, 1.00, y_backstop, 8.00, 1.00, new_text, size=40, bold=True)


def patch_page_20(slide):
    # These localized white boxes cover only translatable teaching text on the white content background.
    # Formulas and orbital graphics remain untouched.
    add_white_box(slide, 1.05, 1.15, 3.85, 2.22)
    add_text(slide, 1.08, 1.20, 3.78, 1.20, '• sp³ 杂化轨道由 s 与 p 波函数的不同加减组合得到。', size=23, bold=False, align=PP_ALIGN.LEFT)

    add_white_box(slide, 4.75, 2.03, 1.15, 0.34)
    add_text(slide, 4.67, 2.03, 1.30, 0.34, '1 个 s 轨道', size=12)

    add_white_box(slide, 5.93, 2.03, 1.25, 0.34)
    add_text(slide, 5.88, 2.03, 1.38, 0.34, '3 个 p 轨道', size=12)

    add_white_box(slide, 8.05, 2.03, 1.45, 0.34)
    add_text(slide, 7.92, 2.03, 1.75, 0.34, '4 个 sp³ 杂化轨道', size=12)

    add_white_box(slide, 6.18, 2.90, 1.55, 0.45)
    add_text(slide, 6.18, 2.90, 1.55, 0.45, '109.5° 键角', size=13)

    add_white_box(slide, 7.02, 4.55, 1.32, 0.42)
    add_text(slide, 7.02, 4.55, 1.32, 0.42, '四面体几何', size=13)

    add_white_box(slide, 5.03, 5.42, 4.55, 0.72)
    add_text(slide, 5.04, 5.42, 4.50, 0.72, '说明：应变会使键角发生畸变；\n例如氨等体系中的键角通常略小。', size=13, align=PP_ALIGN.LEFT)


def set_notes(prs: Presentation, page: int, text: str):
    tf = prs.slides[page - 1].notes_slide.notes_text_frame
    tf.text = text.strip()


def chinese_count(text: str) -> int:
    return len(re.findall(r'[\u4e00-\u9fff]', text))


notes_text = {
8: '''外延生长是指在已有晶体表面继续生长具有确定晶向关系的薄层材料。这里先把后续要讨论的方法放在一张总览中：液相外延利用液体溶液中的过饱和组分在衬底上结晶；气相外延利用气相反应物在高温表面发生反应或分解；分子束外延在超高真空中用定向束流逐层沉积；化学气相沉积和金属有机化学气相沉积则通过气相前驱体输运、分解和表面反应形成薄膜。这些方法都围绕同一个目标：控制材料组分、厚度、界面和晶体质量，为器件结构提供可设计的外延层。

[Sources]
ECE340_L5_S18_Posted.pdf, page 8.''',
9: '''气相外延的核心是把反应气体按控制好的流量送入石英反应腔，使气体沿晶圆表面流动，并在高温区域发生化学反应。晶圆通常放在受热基座上，RF 加热帮助提供表面反应所需的温度。以硅外延为例，SiCl₄ 与 H₂ 可以在表面反应生成 Si 和 HCl，SiH₄ 也可以热分解生成 Si 和 H₂。反应后的副产物和未反应气体沿排气方向带走，所以气流组织、温度分布和晶圆位置都会影响薄膜均匀性。

[Sources]
ECE340_L5_S18_Posted.pdf, page 9.''',
10: '''早期 MOCVD 系统的照片展示了这种工艺在工程实现上的复杂性。金属有机前驱体和氢化物气体需要通过多路气路输送，流量、压力和切换时序都要被精确控制。反应腔负责把气体带到受热衬底表面，使前驱体分解并形成目标化合物半导体薄层。早期系统往往管路密集、操作复杂，但它已经体现出现代 MOCVD 的关键要素：前驱体供给、质量流量控制、反应区温度控制以及废气排放处理。

[Sources]
ECE340_L5_S18_Posted.pdf, page 10.''',
11: '''现代 MOCVD 系统已经发展成适合多片晶圆和重复生产的设备平台。左侧晶圆托盘显示多个晶圆可以在同一反应环境中接受气体输运和表面反应，右侧设备照片体现了气路、反应腔和生产线的系统集成。以 GaAs 生长为例，三甲基镓与砷烷反应生成 GaAs，同时产生甲烷副产物：(CH₃)₃Ga + AsH₃ → GaAs + 3CH₄。通过控制前驱体分压、温度和流场，现代系统能够提高厚度均匀性、组分控制和批量制备能力。

[Sources]
ECE340_L5_S18_Posted.pdf, page 11.''',
12: '''分子束外延依赖超高真空环境和具有方向性的原子或分子束。不同束源提供 Si、Al、Ga、As、Be 等元素，快门控制每一路束流何时到达衬底，从而实现层厚、组分和掺杂的精细控制。图中分子束从束源指向加热衬底，GaAs 衬底决定外延层的晶向关系。右侧显微图中 4×4、GaAs substrate、10 nm 和 <100> 等信息说明了表面重构、衬底材料、尺度和晶向，这些都是理解 MBE 生长质量的重要线索。

[Sources]
ECE340_L5_S18_Posted.pdf, page 12.''',
14: '''生产型 MBE 反应器强调从实验研究走向可重复制备。与研究型系统相比，生产型设备通常具有更大的真空腔体、更完整的多源配置、更稳定的温度和束流控制，以及便于维护和装载的工程结构。它的目标不仅是实现高质量外延层，还要保证不同批次之间的厚度、组分和界面具有一致性。通过这些设备，可以把 MBE 的精确生长能力应用到更接近器件制造的场景中。

[Sources]
ECE340_L5_S18_Posted.pdf, page 14.''',
16: '''固体中的键合类型可以从电子如何重新分布来理解。离子键中，电子主要从一个原子转移到另一个原子，库仑吸引和原子核之间的排斥在平衡距离处达到平衡。共价键中，两个原子通过共享电子形成定向键合，Si、Ge 和 C 都是典型例子。混合离子-共价键出现在电负性差不为零但又没有完全电子转移的情形，电子云发生极化，因此 GaAs、InP、GaN 等化合物半导体具有极性共价特征。金属键中，正离子实浸在离域电子海中，电子不局限于单个键，这解释了金属中电子较容易运动的特点。电负性差 ΔE 从 0 到较大时，键合可以从共价逐渐过渡到极性共价再到离子特征。

[Sources]
ECE340_L5_S18_Posted.pdf, page 16.''',
17: '''这张表把元素的原子轨道电子占据按照 1s、2s、2p、3s、3p 等轨道列出。He core、Ne core 和 Ar core 表示内层电子已经形成相对稳定的芯层结构，后面的价电子决定了元素参与成键的主要方式。C 和 Si 被重点标出，是因为它们都具有四个价电子，能够形成四个方向性的共价键。对于 Si，电子组态可以看作芯层之外的 3s²3p²，这为后面讨论 sp³ 杂化、四面体键合和晶体能带形成打下基础。

[Sources]
ECE340_L5_S18_Posted.pdf, page 17.''',
18: '''硅原子核电荷为 +14，周围共有 14 个电子。内层 1s、2s 和 2p 轨道容纳 10 个芯层电子，它们能量较低、离原子核更近，通常不直接决定半导体成键性质。外层 3s 和 3p 轨道包含 4 个价电子，是形成共价键和能带的关键。右侧能级图把芯层电子、价电子、价轨道、首次激发轨道以及电离或零能级放在同一能量坐标中，帮助我们理解电子从束缚态到更高能态的变化，以及为什么价层电子在固体性质中最重要。

[Sources]
ECE340_L5_S18_Posted.pdf, page 18.''',
19: '''原子轨道的空间形状来自波函数在三维空间中的分布。s 轨道近似球形，没有特定方向；p 轨道具有方向性，常沿 x、y、z 三个方向分布；d 轨道形状更复杂，带有更多节面和取向差异。这些图不是电子绕核运动的经典轨道，而是描述电子出现概率和波函数相位特征的空间图像。理解轨道方向性很重要，因为后续共价键、杂化轨道和晶体中的成键关系都取决于轨道如何相互重叠。

[Sources]
ECE340_L5_S18_Posted.pdf, page 19.''',
20: '''sp³ 杂化可以理解为一个 s 轨道和三个 p 轨道重新组合，形成四个能量等价、空间方向不同的杂化轨道。四条 ψ₁ 到 ψ₄ 的公式都有共同的 1/2 系数，表示每个杂化轨道都由 ψs、ψpx、ψpy 和 ψpz 按相同权重组合而成。它们之间的区别不在于参与的轨道种类，而在于 px、py、pz 各项前面的正负号不同。正负号改变了波函数在空间中的相位组合，因此四个杂化轨道指向不同方向。由于四个组合在数学上对称，它们形成四个彼此等价的 sp³ 轨道。为了使轨道之间的排斥尽可能小，四个方向在空间中排成四面体结构，理想键角约为 109.5°。在真实晶体或分子中，外部应变、邻近原子的影响或分子结构差异会使键角偏离理想值，例如氨等体系的键角会略小。讲解这页时不需要逐字念每条公式，而要强调：s 与 p 的不同加减组合把原本的原子轨道转化为四个指向四面体顶点的成键方向。

[Sources]
ECE340_L5_S18_Posted.pdf, page 20.''',
22: '''能带的形成可以从晶格周期势中的电子波来定性理解。严格来说，需要在周期性势能背景下求解薛定谔方程，得到电子能量 E 与晶体动量波矢 k 之间的关系，也就是 E-k 关系。完整理论超出本课程范围，但定性图像是：周期结构会使电子波发生相干叠加和选择性传播，从而产生允许能量范围和禁带。右侧用光学干涉滤光片和蝴蝶翅膀微结构作类比，是因为它们也依赖周期结构与波长之间的匹配，只是这里对应的是电子波长而不是可见光波长。

[Sources]
ECE340_L5_S18_Posted.pdf, page 22.''',
23: '''从两个氢原子形成氢分子可以直观理解能级分裂。每个 H 原子的 1s 轨道有 2 个可用状态并含有 1 个电子；两个 H 原子靠近后，H₂ 共有 4 个状态和 2 个电子。在线性组合原子轨道方法中，两个 1s 波函数可以相加形成成键轨道，也可以相减形成反键轨道。成键轨道使电子密度更多位于两核之间，降低体系能量；反键轨道使电子密度不集中在两核之间，能量较高。V(r) 图展示了原子间距变化时的相互作用势和能级关系。把这种从两个原子产生能级分裂的思想推广到 N 个原子，原来的离散能级会分裂成大量密集状态，最终形成固体中的能带。

[Sources]
ECE340_L5_S18_Posted.pdf, page 23.''',
24: '''硅原子的电子组态为 1s²2s²2p⁶3s²3p²，其中 1s²2s²2p⁶ 构成芯层电子，3s²3p² 是价电子层。对于 n = 3 的外层，3s 轨道提供 2 个可用状态并已有 2 个电子，3p 轨道提供 6 个可用状态但只有 2 个电子，因此总共有 8 个可用状态和 4 个价电子。当 N 个 Si 原子形成晶体时，这些状态数按 N 倍扩展，总可用状态数为 N×8，总电子数为 N×4。在 0 K 下，较低能量的价带被 N×4 个电子填满，而较高能量的导带为空。这个状态计数说明了离散原子能级如何在晶体中演化为价带、导带和禁带结构。

[Sources]
ECE340_L5_S18_Posted.pdf, page 24.'''
}

META_TERMS = ['rebuilt', '重建', '返修', '监工', 'supervisor', 'ROUND', '裁图', '源页', '原页保留', '原页裁取', '不添加中文贴纸', '白底覆盖', '覆盖英文', '已完成中文化', '冻结', '保持不动', '本页已完成', '待替换', 'Placeholder']

assert BASE.exists(), f'Missing Final Candidate base: {BASE}'
assert SRC.exists(), f'Missing source PDF: {SRC}'
shutil.copyfile(BASE, OUT)

base_slide_xml = slide_xmls(BASE)
prs = Presentation(OUT)
assert len(prs.slides) == 52, f'Expected 52 slides, got {len(prs.slides)}'

# Snapshot frozen notes before changes.
frozen_notes_before = {p: extract_notes(prs, p) for p in NOTES_FROZEN}

# Page 13: delete page number and Chinese caption only.
slide = prs.slides[12]
remove_bottom_page_number(slide, '13', prs)
replace_or_cover_caption_13(slide)

# Page 15: chapter transition title Chinese; remove page number.
slide = prs.slides[14]
remove_bottom_page_number(slide, '15', prs)
replace_big_title(slide, 'Bonding Forces and Bond Types', '固体中的键合力与键合类型', 2.70)

# Page 20: formulas remain under original image; only visible English teaching text is localized.
slide = prs.slides[19]
patch_page_20(slide)

# Page 21: chapter transition title Chinese; remove page number.
slide = prs.slides[20]
remove_bottom_page_number(slide, '21', prs)
replace_big_title(slide, 'Energy Bands', '能带', 3.05)

# Notes changes only on allowed pages.
for page, txt in notes_text.items():
    set_notes(prs, page, txt)

# Verify notes for frozen pages unchanged.
for p in NOTES_FROZEN:
    assert extract_notes(prs, p) == frozen_notes_before[p], f'Frozen notes changed unexpectedly: page {p}'

prs.save(OUT)

# Verify slide XML changes are limited to allowed visual pages.
out_slide_xml = slide_xmls(OUT)
visual_changed = [i for i in range(1, 53) if base_slide_xml[i] != out_slide_xml[i]]
unexpected_visual = [i for i in visual_changed if i not in VISUAL_ALLOWED]
if unexpected_visual:
    raise AssertionError(f'Unexpected visual slide XML changes: {unexpected_visual}; all changed={visual_changed}')

# Notes QA: real per-page check.
prs_check = Presentation(OUT)
notes_rows = []
notes_failed = []
page20_chars = 0
for i in PAGES:
    notes = extract_notes(prs_check, i)
    cc = chinese_count(notes)
    has_sources = '[Sources]' in notes
    meta_hits = [term for term in META_TERMS if term in notes]
    body = notes.split('[Sources]')[0].strip() if '[Sources]' in notes else notes.strip()
    body_cc = chinese_count(body)
    pure_sources = body_cc < 30
    if i == 20:
        page20_chars = body_cc
    ok = bool(notes.strip()) and has_sources and not meta_hits and not pure_sources and (i != 20 or body_cc >= 400)
    notes_rows.append((i, cc, body_cc, has_sources, meta_hits, ok))
    if not ok:
        notes_failed.append(i)
if notes_failed:
    raise AssertionError(f'Notes QA failed pages: {notes_failed}; rows={notes_rows}')

# Export PDF and render evidence.
subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', str(PDF_DIR), str(OUT)], check=True)
pdf_path = PDF_DIR / (OUT.stem + '.pdf')
assert pdf_path.exists(), f'Missing exported PDF: {pdf_path}'
newdoc = fitz.open(pdf_path)
origdoc = fitz.open(SRC)
assert newdoc.page_count == 52, f'Final Candidate R1 PDF pages: {newdoc.page_count}'
assert origdoc.page_count == 52, f'Original PDF pages: {origdoc.page_count}'

render_paths = []
comparison_paths = []
for page in PAGES:
    pix = newdoc[page - 1].get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
    png_path = RENDER / f'page_{page:02d}.png'
    pix.save(png_path)
    render_paths.append(png_path)

    opix = origdoc[page - 1].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
    npix = newdoc[page - 1].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
    orig = Image.frombytes('RGB', [opix.width, opix.height], opix.samples)
    new = Image.frombytes('RGB', [npix.width, npix.height], npix.samples)
    canvas = Image.new('RGB', (orig.width + new.width + 70, max(orig.height, new.height) + 80), 'white')
    d = ImageDraw.Draw(canvas)
    d.text((10, 10), f'Original PDF page {page}', fill=(0, 0, 0))
    d.text((orig.width + 60, 10), f'Final Candidate R1 page {page}', fill=(0, 0, 0))
    canvas.paste(orig, (10, 55))
    canvas.paste(new, (orig.width + 60, 55))
    comp_path = COMP / f'page_{page:02d}_original_pdf_vs_final_candidate_r1.jpg'
    canvas.save(comp_path, quality=92)
    comparison_paths.append(comp_path)

# Contact sheet.
thumbs = []
for page, path in zip(PAGES, render_paths):
    im = Image.open(path).convert('RGB')
    im.thumbnail((360, 235))
    tile = Image.new('RGB', (390, 280), 'white')
    d = ImageDraw.Draw(tile)
    d.text((10, 8), f'page {page:02d}', fill=(0, 0, 0))
    tile.paste(im, ((390 - im.width) // 2, 38))
    thumbs.append(tile)
cols, rows = 4, 5
sheet = Image.new('RGB', (cols * 390, rows * 280), 'white')
for idx, tile in enumerate(thumbs):
    sheet.paste(tile, ((idx % cols) * 390, (idx // cols) * 280))
contact_path = EV / 'contact_sheet_final_qa_r1_pages_08_24.jpg'
sheet.save(contact_path, quality=92)

out_hash = sha256(OUT)
notes_table = '\n'.join(
    f'| {page} | {cc} | {body_cc} | {"yes" if has_sources else "no"} | {", ".join(meta_hits) if meta_hits else "none"} | {"passed" if ok else "failed"} |'
    for page, cc, body_cc, has_sources, meta_hits, ok in notes_rows
)
report = f'''# ECE340 L5 Final QA R1 Build Report（第 8–24 页）

- Supervisor feedback source: `l5/SUPERVISOR_FINAL_QA_VISUAL_REVIEW_FEEDBACK_ROUND1.md`, commit `ebf1fb43a7b6346e45621d11e8a47bf0a300ab81`.
- 基准 Final Candidate：`{BASE}`
- 新 Final Candidate R1：`{OUT}`
- Final Candidate R1 SHA-256：`{out_hash}`
- 52 页页数确认：passed（PPT: {len(prs_check.slides)} slides; exported PDF: {newdoc.page_count} pages）
- 实际视觉修改页面：第 13、15、20、21 页
- 实际 Notes 修改页面：第 8、9、10、11、12、14、16、17、18、19、20、22、23、24 页
- 第 13、15、21 页 notes 冻结确认：passed
- 第 8、9、10、11、12、14、16、17、18、19、22、23、24 页视觉冻结确认：passed
- 第 1–7、25–52 页视觉冻结确认：passed
- Slide XML changed pages: {visual_changed}
- Unexpected visual slide XML changes: {unexpected_visual}

## 17 张 PNG 路径

''' + '\n'.join(f'- `{p}`' for p in render_paths) + f'''

## 17 张 Original PDF vs Final Candidate R1 对照图路径

''' + '\n'.join(f'- `{p}`' for p in comparison_paths) + f'''

## Contact sheet

- `{contact_path}`

## Notes QA

Notes QA: passed

| page | Chinese chars total | Chinese chars before Sources | has [Sources] | meta-term hits | status |
|---:|---:|---:|:---:|:---|:---:|
{notes_table}

- 第 20 页讲稿有效汉字数：{page20_chars}

## Visual regression

- Evidence generated: 17 PNG + 17 Original PDF vs Final Candidate R1 comparison images + contact sheet.
- Actual 17-page visual inspection by model after artifact download: pending at build time.

## Microsoft PowerPoint actual open check

- Microsoft PowerPoint actual open check: not performed / pending

## Supervisor status

- Supervisor final acceptance: pending
'''
REPORT.write_text(report, encoding='utf-8')

subprocess.run(['git', 'add', str(OUT), str(REPORT), str(EV)], check=True)
print('Final QA R1 package complete')
print('PPT:', OUT)
print('Rendered:', RENDER)
print('Comparison:', COMP)
print('Contact:', contact_path)
print('Report:', REPORT)
print('Page20 Chinese chars:', page20_chars)
print('Visual changed pages:', visual_changed)
