from pathlib import Path
import shutil
import zipfile
import hashlib
import subprocess
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_最终候选版_R1_第8-24页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_最终候选版_R2_第8-24页.pptx'
EV = ROOT / 'final_qa_r2_page18'
RENDER = EV / 'rendered'
COMP = EV / 'comparison'
PDF_DIR = EV / 'final_pdf'
REPORT = ROOT / 'BUILD_REPORT_FINAL_QA_R2_PAGE18.md'
PAGE = 18
for p in (EV, RENDER, COMP, PDF_DIR):
    p.mkdir(parents=True, exist_ok=True)


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def slide_xmls(pptx: Path):
    data = {}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1, 53):
            data[i] = z.read(f'ppt/slides/slide{i}.xml')
    return data


def notes_xmls(pptx: Path):
    data = {}
    with zipfile.ZipFile(pptx) as z:
        for name in z.namelist():
            if name.startswith('ppt/notesSlides/notesSlide') and name.endswith('.xml'):
                data[name] = z.read(name)
    return data


def add_white_box(slide, x, y, w, h):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shp.line.color.rgb = RGBColor(255, 255, 255)
    return shp


def add_text(slide, x, y, w, h, text, size=14):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = 'Noto Sans CJK SC'
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(0, 0, 0)
    return box


def patch_page18_caption(prs: Presentation):
    slide = prs.slides[PAGE - 1]
    # The existing caption is black text printed on the white figure background.
    # This white cover stays inside the white caption strip and does not touch the science diagrams.
    add_white_box(slide, 1.48, 5.03, 7.15, 1.05)
    add_text(
        slide,
        1.58,
        5.24,
        6.95,
        0.58,
        '图 2.8　Si 原子的电子结构与能级示意：（a）轨道模型显示 10 个芯层电子（n = 1、2）和 4 个价电子（n = 3）；（b）示意给出原子核库仑势中的电子能级。',
        size=14,
    )


def render_page18():
    subprocess.run([
        'libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', str(PDF_DIR), str(OUT)
    ], check=True)
    pdf = PDF_DIR / (OUT.stem + '.pdf')
    doc = fitz.open(str(pdf))
    page = doc[PAGE - 1]
    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    target = RENDER / 'page_18.png'
    pix.save(str(target))
    return pdf, target


def make_comparison(page_png: Path):
    srcdoc = fitz.open(str(SRC))
    spage = srcdoc[PAGE - 1]
    spix = spage.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    orig = EV / 'original_pdf_page_18.png'
    spix.save(str(orig))
    left = Image.open(orig).convert('RGB')
    right = Image.open(page_png).convert('RGB')
    h = max(left.height, right.height)
    lw = int(left.width * h / left.height)
    rw = int(right.width * h / right.height)
    left = left.resize((lw, h))
    right = right.resize((rw, h))
    margin = 40
    header = 72
    out = Image.new('RGB', (lw + rw + margin * 3, h + header + margin), 'white')
    out.paste(left, (margin, header))
    out.paste(right, (margin * 2 + lw, header))
    draw = ImageDraw.Draw(out)
    draw.text((margin, 24), 'Original PDF page 18', fill=(0, 0, 0))
    draw.text((margin * 2 + lw, 24), 'Final Candidate R2 page 18', fill=(0, 0, 0))
    path = COMP / 'page_18_original_pdf_vs_final_candidate_r2.jpg'
    out.save(path, quality=92)
    return path


def main():
    if not BASE.exists():
        raise FileNotFoundError(BASE)
    before_slide = slide_xmls(BASE)
    before_notes = notes_xmls(BASE)
    prs = Presentation(str(BASE))
    if len(prs.slides) != 52:
        raise AssertionError(f'Expected 52 slides, found {len(prs.slides)}')
    patch_page18_caption(prs)
    prs.save(str(OUT))
    after_slide = slide_xmls(OUT)
    changed_slides = [i for i in range(1, 53) if before_slide[i] != after_slide[i]]
    if changed_slides != [PAGE]:
        raise AssertionError(f'Unexpected changed slides: {changed_slides}')
    after_notes = notes_xmls(OUT)
    if before_notes != after_notes:
        raise AssertionError('Notes changed; R2 must not modify notes.')
    pdf, page_png = render_page18()
    comp = make_comparison(page_png)
    report = f'''# ECE340 L5 Final QA R2 Page 18 Caption Build Report

- Supervisor feedback source: `l5/SUPERVISOR_FINAL_QA_VISUAL_REVIEW_FEEDBACK_ROUND2.md`, commit `28e772661c359167ed8baffd27f95cd82b8d27fb`.
- 基准 Final Candidate R1：`{BASE}`
- 新 Final Candidate R2：`{OUT}`
- Final Candidate R2 SHA-256：`{sha256(OUT)}`
- 52 页页数确认：passed
- 本轮修改页面：第 18 页
- 本轮修改区域：第 18 页底部英文 Figure 2.8 图注区域
- 科学图内部修改：无
- Notes 修改：无
- Slide XML changed pages: {changed_slides}
- 第 8–17、19–24 页冻结确认：passed
- 第 1–7、25–52 页冻结确认：passed

## 渲染证据

- 第 18 页高清 PNG：`{page_png}`
- 第 18 页 Original PDF vs Final Candidate R2 对照图：`{comp}`

## Worker self-check

- Only slide 18 changed: passed
- Notes unchanged: passed
- Page count: 52

## Supervisor status

- Supervisor final acceptance: pending
'''
    REPORT.write_text(report, encoding='utf-8')
    subprocess.run(['git', 'add', str(OUT), str(page_png), str(comp), str(REPORT)], check=True)


if __name__ == '__main__':
    main()
