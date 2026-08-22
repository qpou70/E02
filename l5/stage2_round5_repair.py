from pathlib import Path
import zipfile
import hashlib
import subprocess
import fitz
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path('l5')
BASE = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND4_第16页.pptx'
SRC = ROOT / 'stage1_source_reference/ECE340_L5_S18_Posted.pdf'
OUT = ROOT / 'ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND5_第16页.pptx'
ASSET = ROOT / 'stage2_round5_assets'
EV = ROOT / 'stage2_visual_review_round5'
RENDER = EV / 'rendered'
COMP = EV / 'comparison'
REPORT = ROOT / 'BUILD_REPORT_STAGE2_ROUND5.md'
TARGET = [16]
for p in [ASSET, RENDER, COMP]:
    p.mkdir(parents=True, exist_ok=True)

NAVY = RGBColor(30, 50, 82)
BLACK = RGBColor(28, 28, 30)
WHITE = RGBColor(255, 255, 255)
MID = RGBColor(204, 216, 230)


def slide_xmls(pptx):
    data = {}
    with zipfile.ZipFile(pptx) as z:
        for i in range(1, 53):
            data[i] = z.read(f'ppt/slides/slide{i}.xml')
    return data


def sha256(path):
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


def inches(value):
    return value / 914400


def remove_shapes_in_region(slide, x0, y0, x1, y1):
    removed = 0
    for sh in list(slide.shapes):
        sx0 = inches(sh.left)
        sy0 = inches(sh.top)
        sx1 = sx0 + inches(sh.width)
        sy1 = sy0 + inches(sh.height)
        # Remove only shapes whose bounding box overlaps the accepted metal-bonding visual zone.
        if sx1 > x0 and sx0 < x1 and sy1 > y0 and sy0 < y1:
            sh.element.getparent().remove(sh.element)
            removed += 1
    return removed


def textbox(slide, x, y, w, h, text, size=10, bold=False, color=BLACK, align=PP_ALIGN.CENTER, fill=None, line=None, margin=2):
    sh = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(0.8)
    tf = sh.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(margin)
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = 'Noto Sans CJK SC'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return sh


def render_crop(doc, page_no, coords, name, dpi=360):
    pix = doc[page_no - 1].get_pixmap(matrix=fitz.Matrix(dpi/72, dpi/72), clip=fitz.Rect(*coords), alpha=False)
    out = ASSET / name
    pix.save(out)
    return out


def add_pic(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    k = min(w / iw, h / ih)
    pw, ph = iw * k, ih * k
    return slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))

prs = Presentation(BASE)
doc = fitz.open(SRC)
assert len(prs.slides) == 52 and doc.page_count == 52
base_xml = slide_xmls(BASE)
slide = prs.slides[15]

# Only replace the right-lower metallic-bonding scientific figure area.
# Keep the accepted left-side four text boxes, upper lattice figure, the Chinese shared-electron note,
# and the three Delta-E electron-cloud labels/figures intact.
removed = remove_shapes_in_region(slide, 3.60, 4.46, 9.55, 6.55)

# Crop only the true metallic-bonding science region from original PDF page 16:
# left positive cores + orange delocalised electrons, centre red arrow, right electron cloud + positive cores.
# The crop excludes original English title, original English note, long caption, and external whitespace.
metal_clean = render_crop(doc, 16, (405, 462, 670, 503), 'round5_p16_metallic_science_clean.png', 380)

textbox(slide, 3.86, 4.50, 5.25, 0.24, '金属键：正离子实 + 离域电子海', 9.2, True, NAVY, PP_ALIGN.CENTER)
add_pic(slide, metal_clean, 3.78, 4.83, 5.58, 0.96)
textbox(slide, 4.42, 5.93, 2.80, 0.34, '离域电子海', 10.0, True, BLACK, PP_ALIGN.CENTER)
slide.notes_slide.notes_text_frame.text = '[Sources]\nECE340_L5_S18_Posted.pdf, page 16. ROUND5 only replaces the clean metallic-bonding science image in the lower-right area; all other accepted slide-16 content and all other slides remain frozen.'

# Text checks on page 16.
visible_text = '\n'.join((getattr(sh, 'text', '') or '') for sh in slide.shapes)
for required in ['ΔE = 0（共价键）', 'ΔE 中等（极性共价键）', 'ΔE 较大（离子键）', '每个键由两个电子共享', '金属键：正离子实 + 离域电子海', '离域电子海']:
    assert required in visible_text, f'missing accepted or required text: {required}'
for forbidden in ['Metallic Bonding\n', 'Swarm of delocalised electrons', 'Two electrons per bond', 'Placeholder', '待替换', '已清理', '已重建', '已中文化', '去除英文']:
    assert forbidden not in visible_text, f'forbidden visible text remains: {forbidden}'

prs.save(OUT)
post_xml = slide_xmls(OUT)
changed = [i for i in range(1, 53) if post_xml[i] != base_xml[i]]
if changed != TARGET:
    raise AssertionError(f'Only page 16 may change; actual changed pages: {changed}')
pre_hash = sha256(OUT)

# Export PDF and render evidence from GitHub build.
pdf_dir = EV / 'new_pdf'
pdf_dir.mkdir(parents=True, exist_ok=True)
subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', str(pdf_dir), str(OUT)], check=True)
pdf_path = pdf_dir / (OUT.stem + '.pdf')
newdoc = fitz.open(pdf_path)
assert newdoc.page_count == 52
pix = newdoc[15].get_pixmap(matrix=fitz.Matrix(2.2, 2.2), alpha=False)
page_png = RENDER / 'page_16.png'
pix.save(page_png)

# Original PDF vs New Page comparison.
opix = doc[15].get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
orig = Image.frombytes('RGB', [opix.width, opix.height], opix.samples)
new = Image.open(page_png).convert('RGB')
canvas = Image.new('RGB', (orig.width + new.width + 60, max(orig.height, new.height) + 70), 'white')
d = ImageDraw.Draw(canvas)
d.text((10, 10), 'Original PDF page 16', fill=(0, 0, 0))
d.text((orig.width + 50, 10), 'New Stage 2 ROUND5 page 16', fill=(0, 0, 0))
canvas.paste(orig, (10, 50))
canvas.paste(new, (orig.width + 50, 50))
comp_path = COMP / 'page_16_original_pdf_vs_new_round5.jpg'
canvas.save(comp_path, quality=92)

# Single-page contact sheet.
thumb = new.copy()
thumb.thumbnail((620, 460))
cs = Image.new('RGB', (680, 520), 'white')
d = ImageDraw.Draw(cs)
d.text((10, 10), 'page 16', fill=(0, 0, 0))
cs.paste(thumb, ((680 - thumb.width) // 2, 45))
contact_path = EV / 'contact_sheet_stage2_round5_page_16.jpg'
cs.save(contact_path, quality=92)
post_hash = sha256(OUT)

report = f'''# ECE340 L5 第二阶段视觉返修 ROUND5 Build Report

- 返修依据：`l5/SUPERVISOR_STAGE2_VISUAL_REVIEW_FEEDBACK_ROUND5.md`，反馈提交 `7a497e8a0bafe6bd17d428c252d4bd72f06580fd`。
- 基准 PPT 路径：`l5/ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND4_第16页.pptx`
- 新 PPT 路径：`l5/ECE340_L5_S18_Posted_中文忠实重建_第二阶段视觉返修ROUND5_第16页.pptx`
- 本轮只修改第 16 页右下金属键图区。
- 其他 slide 冻结：第 8–15、17–24 页；第 1–7、25–52 页。
- Slide XML 检查：仅第 16 页与 ROUND4 基准不同。
- 删除对象数量：{removed} 个，仅限右下金属键区域重建。

## 第 16 页右下金属键图区处理

- 从 Original PDF 第 16 页重新干净裁取金属键科学图的实际科学区域。
- 裁取内容只包括：左侧正离子实与小橙色离域电子、中央红色箭头、右侧电子云中的正离子实。
- 不裁入 `Metallic Bonding`、`Swarm of delocalised electrons`、英文长段说明或原图外部空白。
- 保留中文标题 `金属键：正离子实 + 离域电子海` 与中文说明 `离域电子海`。
- 保持三幅 ΔE 标签和 `每个键由两个电子共享` 不变。

## 渲染证据

- 渲染 PNG 路径：`l5/stage2_visual_review_round5/rendered/page_16.png`
- comparison 路径：`l5/stage2_visual_review_round5/comparison/page_16_original_pdf_vs_new_round5.jpg`
- contact sheet 路径：`l5/stage2_visual_review_round5/contact_sheet_stage2_round5_page_16.jpg`

## 自检记录

- Worker self-check 状态：passed。
- Supervisor visual acceptance: pending。
- PPT SHA-256（渲染前）：`{pre_hash}`
- PPT SHA-256（渲染后）：`{post_hash}`
'''
REPORT.write_text(report, encoding='utf-8')

# Stage all generated outputs.
subprocess.run(['git', 'add', str(OUT), str(REPORT), str(EV)], check=True)
print('ROUND5 page 16 metallic-bonding cleanup build complete')
print('PPT:', OUT)
print('PNG:', page_png)
print('Comparison:', comp_path)
print('Contact:', contact_path)
