from pathlib import Path
import hashlib, shutil, subprocess, zipfile, re, textwrap
from PIL import Image, ImageDraw, ImageFont
import fitz
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
BASE_PPT = ROOT / 'l5/ECE340_L5_S18_Posted_中文忠实重建_最终候选版_R2_第8-24页.pptx'
OUT_PPT = ROOT / 'l5/ECE340_L5_S18_Posted_中文忠实重建_最终候选版_R3_第8-24页.pptx'
WORK = ROOT / 'l5/final_qa_image_rebuild_20_22_23'
GEN_DIR = WORK / 'generated'
RENDER_DIR = WORK / 'rendered'
COMP_DIR = WORK / 'comparison'
PDF_DIR = WORK / 'pdf'
REPORT = ROOT / 'l5/BUILD_REPORT_FINAL_QA_IMAGE_REBUILD_20_22_23.md'
CONTACT = WORK / 'contact_sheet_pages_20_22_23.jpg'
SOURCE_DIR = ROOT / 'l5/stage1_source_reference'
PAGES = [20,22,23]
for d in [GEN_DIR, RENDER_DIR, COMP_DIR, PDF_DIR]: d.mkdir(parents=True, exist_ok=True)

FONT_REG = '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc'
FONT_BOLD = '/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc'
FONT_SERIF = '/usr/share/fonts/truetype/dejavu/DejaVuSerif.ttf'

def font(size, bold=False, serif=False):
    if serif:
        return ImageFont.truetype(FONT_SERIF, size)
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REG, size)

def rect(draw, box, fill='white'):
    draw.rectangle([int(x) for x in box], fill=fill)

def draw_wrapped(draw, xy, text, fnt, fill='black', max_width=600, line_spacing=1.18):
    x,y = xy
    lines=[]
    for para in text.split('\n'):
        cur=''
        for ch in para:
            test=cur+ch
            if draw.textbbox((0,0), test, font=fnt)[2] <= max_width or not cur:
                cur=test
            else:
                lines.append(cur); cur=ch
        if cur: lines.append(cur)
    step=int(fnt.size*line_spacing)
    for i,line in enumerate(lines):
        draw.text((x,y+i*step), line, font=fnt, fill=fill)
    return y+len(lines)*step

def make_page20():
    im = Image.open(SOURCE_DIR/'original_pdf_page_20.png').convert('RGB')
    W,H=im.size; d=ImageDraw.Draw(im)
    blue=(9,37,150); orange=(255,128,0)
    rect(d, (0.045*W,0.09*H,0.955*W,0.205*H), blue)
    rect(d, (0.045*W,0.205*H,0.955*W,0.225*H), orange)
    d.text((0.075*W,0.115*H), 'sp³ 杂化', font=font(int(0.058*H), True), fill='white')
    rect(d, (0.10*W,0.27*H,0.47*W,0.53*H), 'white')
    d.text((0.105*W,0.285*H), '•', font=font(int(0.047*H), True), fill='black')
    draw_wrapped(d, (0.145*W,0.285*H), 'sp³ 杂化轨道由 s 与 p 波函数的不同加减组合得到。', font(int(0.036*H), True), max_width=int(0.32*W), line_spacing=1.25)
    rect(d, (0.47*W,0.39*H,0.56*W,0.43*H),'white')
    rect(d, (0.595*W,0.39*H,0.70*W,0.43*H),'white')
    rect(d, (0.82*W,0.39*H,0.955*W,0.43*H),'white')
    d.text((0.475*W,0.395*H),'1 个 s 轨道',font=font(int(0.022*H), True),fill='black')
    d.text((0.600*W,0.395*H),'3 个 p 轨道',font=font(int(0.022*H), True),fill='black')
    d.text((0.825*W,0.395*H),'4 个 sp³ 杂化轨道',font=font(int(0.021*H), True),fill='black')
    rect(d, (0.64*W,0.49*H,0.78*W,0.57*H),'white')
    d.text((0.655*W,0.505*H),'109.5° 键角',font=font(int(0.026*H), True),fill='black')
    rect(d, (0.70*W,0.72*H,0.84*W,0.79*H),'white')
    d.text((0.72*W,0.735*H),'四面体几何',font=font(int(0.026*H), True),fill='black')
    rect(d,(0.49*W,0.805*H,0.95*W,0.965*H),'white')
    d.text((0.51*W,0.825*H),'说明：',font=font(int(0.024*H), True),fill='black')
    draw_wrapped(d,(0.51*W,0.86*H),'应变会使键角发生畸变；\n例如氨等体系中的键角通常略小。',font(int(0.022*H), True),max_width=int(0.40*W),line_spacing=1.25)
    im.save(GEN_DIR/'page_20_cn.png')

def make_page22():
    im = Image.open(SOURCE_DIR/'original_pdf_page_22.png').convert('RGB')
    W,H=im.size; d=ImageDraw.Draw(im)
    blue=(9,37,150); orange=(255,128,0)
    rect(d,(0.045*W,0.09*H,0.955*W,0.205*H),blue); rect(d,(0.045*W,0.205*H,0.955*W,0.225*H),orange)
    d.text((0.075*W,0.115*H),'周期势与 E-k 关系',font=font(int(0.052*H),True),fill='white')
    rect(d,(0.08*W,0.245*H,0.52*W,0.88*H),'white')
    bullets=['Streetman 教材和本课程提供的是能带形成的定性理解。','能带形成的详细理论超出本课程范围。','在晶格周期势中求解薛定谔方程会产生能带。','求解结果给出电子能量 E 与晶体中动量波矢 k 的关系，即 E-k 关系。','这可类比干涉滤光片或蝴蝶翅膀衍射结构，只是作用尺度对应电子波长。']
    y=0.285*H
    for b in bullets:
        d.text((0.105*W,y),'•',font=font(int(0.034*H),True),fill='black')
        y=draw_wrapped(d,(0.145*W,y+0.002*H),b,font(int(0.026*H),False),max_width=int(0.35*W),line_spacing=1.18)+0.035*H
    rect(d,(0.90*W,0.88*H,0.96*W,0.95*H),'white')
    im.save(GEN_DIR/'page_22_cn.png')

def make_page23():
    im = Image.open(SOURCE_DIR/'original_pdf_page_23.png').convert('RGB')
    W,H=im.size; d=ImageDraw.Draw(im)
    blue=(9,37,150); orange=(255,128,0); cyan=(0,160,205)
    rect(d,(0.045*W,0.09*H,0.955*W,0.205*H),blue); rect(d,(0.045*W,0.205*H,0.955*W,0.225*H),orange)
    d.text((0.07*W,0.115*H),'氢分子的成键与反键轨道',font=font(int(0.045*H),True),fill='white')
    rect(d,(0.055*W,0.235*H,0.34*W,0.38*H),'white')
    lines=['氢原子：1s¹','氢原子 1：2 个状态，1 个电子','氢原子 2：2 个状态，1 个电子','H₂：4 个状态，2 个电子']
    y=0.245*H
    for i,line in enumerate(lines):
        d.text((0.065*W,y),line,font=font(int(0.024*H),True if i==0 else False),fill='black')
        y += 0.032*H
    rect(d,(0.63*W,0.25*H,0.90*W,0.36*H),'white')
    d.text((0.735*W,0.265*H),'电子不位于两个氢原子之间',font=font(int(0.021*H),False),fill='black')
    rect(d,(0.78*W,0.36*H,0.95*W,0.49*H),'white')
    d.text((0.83*W,0.365*H),'较高能量',font=font(int(0.033*H),True),fill='black')
    d.text((0.83*W,0.445*H),'较低能量',font=font(int(0.033*H),True),fill='black')
    d.text((0.70*W,0.515*H),'电子位于两个氢原子之间',font=font(int(0.021*H),False),fill='black')
    rect(d,(0.59*W,0.315*H,0.72*W,0.36*H),'white')
    d.text((0.60*W,0.325*H),'反键轨道',font=font(int(0.016*H),True),fill=cyan)
    rect(d,(0.60*W,0.45*H,0.72*W,0.49*H),'white')
    d.text((0.61*W,0.455*H),'成键轨道',font=font(int(0.016*H),True),fill=cyan)
    rect(d,(0.395*W,0.40*H,0.51*W,0.44*H),'white')
    d.text((0.415*W,0.405*H),'原子轨道',font=font(int(0.018*H),True),fill=cyan)
    rect(d,(0.24*W,0.55*H,0.32*W,0.66*H),'white')
    d.text((0.245*W,0.565*H),'反键',font=font(int(0.016*H),False),fill='black')
    d.text((0.245*W,0.635*H),'成键',font=font(int(0.016*H),False),fill='black')
    rect(d,(0.48*W,0.72*H,0.58*W,0.76*H),'white')
    rect(d,(0.61*W,0.535*H,0.78*W,0.61*H),'white')
    d.text((0.625*W,0.545*H),'反键能级',font=font(int(0.017*H),True),fill=cyan)
    d.text((0.625*W,0.595*H),'成键能级',font=font(int(0.017*H),True),fill=cyan)
    rect(d,(0.10*W,0.755*H,0.94*W,0.965*H),'white')
    paragraph='线性组合原子轨道（LCAO）：两个原子靠近时，原子轨道线性组合形成两个不同的“正常”模式——较高能量的反键轨道和较低能量的成键轨道。成键态中电子概率密度在两核之间较高，从而降低成键能级并增强体系内聚；若从两个原子推广到 N 个原子，则会形成 N 个 LCAO 以及 N 个彼此接近的能级，最终形成能带。'
    draw_wrapped(d,(0.14*W,0.785*H),paragraph,font(int(0.018*H),False),max_width=int(0.76*W),line_spacing=1.28)
    im.save(GEN_DIR/'page_23_cn.png')

for f in [make_page20, make_page22, make_page23]: f()

shutil.copy2(BASE_PPT, OUT_PPT)
prs = Presentation(str(OUT_PPT))
assert len(prs.slides) == 52, f'Expected 52 slides, got {len(prs.slides)}'
slide_w, slide_h = prs.slide_width, prs.slide_height
for p in PAGES:
    slide = prs.slides[p-1]
    spTree = slide.shapes._spTree
    for shp in list(slide.shapes):
        spTree.remove(shp._element)
    slide.shapes.add_picture(str(GEN_DIR / f'page_{p:02d}_cn.png'), 0, 0, width=slide_w, height=slide_h)
prs.save(str(OUT_PPT))

def notes_text(ppt_path, page):
    prs = Presentation(str(ppt_path))
    slide = prs.slides[page-1]
    if not slide.has_notes_slide: return ''
    return '\n'.join(getattr(s,'text','') for s in slide.notes_slide.shapes)
for p in PAGES:
    if notes_text(BASE_PPT,p) != notes_text(OUT_PPT,p):
        raise AssertionError(f'Notes changed unexpectedly on slide {p}')

def slide_xml_map(ppt_path):
    out={}
    with zipfile.ZipFile(ppt_path,'r') as z:
        for name in z.namelist():
            m=re.match(r'ppt/slides/slide(\d+)\.xml$',name)
            if m: out[int(m.group(1))]=z.read(name)
    return out
base_xml=slide_xml_map(BASE_PPT); out_xml=slide_xml_map(OUT_PPT)
changed=[i for i in range(1,53) if base_xml.get(i)!=out_xml.get(i)]
if changed != PAGES:
    raise AssertionError(f'Unexpected changed slide XMLs: {changed}')

subprocess.run(['soffice','--headless','--convert-to','pdf','--outdir',str(PDF_DIR),str(OUT_PPT)],check=True)
pdf_candidates=sorted(PDF_DIR.glob('*.pdf'),key=lambda p:p.stat().st_mtime,reverse=True)
if not pdf_candidates: raise FileNotFoundError('LibreOffice did not create a PDF')
pdf_path=pdf_candidates[0]
doc=fitz.open(str(pdf_path))
for p in PAGES:
    pix=doc[p-1].get_pixmap(matrix=fitz.Matrix(2,2),alpha=False)
    pix.save(str(RENDER_DIR/f'page_{p:02d}.png'))
doc.close()

def fit_to_height(img, height):
    return img.resize((int(img.width*height/img.height), height), Image.LANCZOS)
for p in PAGES:
    original=Image.open(SOURCE_DIR/f'original_pdf_page_{p:02d}.png').convert('RGB')
    new=Image.open(RENDER_DIR/f'page_{p:02d}.png').convert('RGB')
    target_h=1000; o=fit_to_height(original,target_h); n=fit_to_height(new,target_h)
    canvas=Image.new('RGB',(o.width+n.width+50,target_h+70),'white')
    canvas.paste(o,(0,70)); canvas.paste(n,(o.width+50,70))
    d=ImageDraw.Draw(canvas); d.text((10,15),f'Original PDF page {p}',fill='black'); d.text((o.width+60,15),f'Final Candidate R3 page {p}',fill='black')
    canvas.save(COMP_DIR/f'page_{p:02d}_original_pdf_vs_final_candidate_r3.jpg',quality=95)
cards=[]
for p in PAGES:
    im=Image.open(RENDER_DIR/f'page_{p:02d}.png').convert('RGB'); im.thumbnail((520,390),Image.LANCZOS)
    card=Image.new('RGB',(540,440),'white'); card.paste(im,((540-im.width)//2,40)); ImageDraw.Draw(card).text((10,10),f'Page {p} - R3 rendered',fill='black')
    cards.append(card)
contact=Image.new('RGB',(540*3,440),'white')
for i,card in enumerate(cards): contact.paste(card,(i*540,0))
contact.save(CONTACT,quality=95)

def sha256(path):
    h=hashlib.sha256()
    with open(path,'rb') as f:
        for c in iter(lambda:f.read(1048576),b''): h.update(c)
    return h.hexdigest()
report=f'''# ECE340 L5 Final QA Image Rebuild 20/22/23

## Scope
- Base PPT: `{BASE_PPT.relative_to(ROOT)}`
- New PPT: `{OUT_PPT.relative_to(ROOT)}`
- Modified visual pages only: 20, 22, 23
- Notes modified: no
- Frozen pages: all pages except 20, 22, 23
- Strategy: full-page Chinese images assembled back into PPT as single full-page pictures.
- Supervisor visual acceptance: pending

## Generated full-page Chinese images
- `{(GEN_DIR/'page_20_cn.png').relative_to(ROOT)}`
- `{(GEN_DIR/'page_22_cn.png').relative_to(ROOT)}`
- `{(GEN_DIR/'page_23_cn.png').relative_to(ROOT)}`

## Rendered evidence
- `{(RENDER_DIR/'page_20.png').relative_to(ROOT)}`
- `{(RENDER_DIR/'page_22.png').relative_to(ROOT)}`
- `{(RENDER_DIR/'page_23.png').relative_to(ROOT)}`

## Original PDF vs R3 comparisons
- `{(COMP_DIR/'page_20_original_pdf_vs_final_candidate_r3.jpg').relative_to(ROOT)}`
- `{(COMP_DIR/'page_22_original_pdf_vs_final_candidate_r3.jpg').relative_to(ROOT)}`
- `{(COMP_DIR/'page_23_original_pdf_vs_final_candidate_r3.jpg').relative_to(ROOT)}`

## Contact sheet
- `{CONTACT.relative_to(ROOT)}`

## Integrity checks
- Slide count: 52
- Changed slide XMLs: {changed}
- Expected changed slide XMLs: {PAGES}
- Notes unchanged on pages 20, 22, 23: yes
- New PPT SHA-256: `{sha256(OUT_PPT)}`

## Worker status
第20、22、23页整页中文图已生成并装配回 PPT，视觉证据已提交，等待 supervisor 检查。
'''
REPORT.write_text(report,encoding='utf-8')
subprocess.run(['git','add',str(OUT_PPT),str(WORK),str(REPORT)],check=True)
print(report)
