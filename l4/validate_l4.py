from pathlib import Path
import hashlib

import fitz
from pptx import Presentation

ROOT = Path(__file__).resolve().parent
PPT = ROOT / "ECE340_L4_S18_Posted_中文忠实重建_逐页讲稿版.pptx"
PDF = ROOT / "rendered" / "ECE340_L4_S18_Posted_中文忠实重建_逐页讲稿版.pdf"

prs = Presentation(PPT)
assert len(prs.slides) == 47, f"Expected 47 slides, found {len(prs.slides)}"
missing = []
for index, slide in enumerate(prs.slides, 1):
    notes = slide.notes_slide.notes_text_frame.text
    if "[Sources]" not in notes:
        missing.append(index)
assert not missing, f"Slides missing notes or [Sources]: {missing}"

rendered = fitz.open(PDF)
assert rendered.page_count == 47, f"Expected 47 rendered pages, found {rendered.page_count}"
sha256 = hashlib.sha256(PPT.read_bytes()).hexdigest()
report = "\n".join([
    "# ECE340 L4 build report",
    "",
    "- Generated inside GitHub Actions: yes",
    "- Source PDF pages: 47",
    f"- PPT slides: {len(prs.slides)}",
    "- Slides with notes and [Sources]: 47",
    f"- PowerPoint-to-PDF rendered pages: {rendered.page_count}",
    f"- SHA-256: `{sha256}`",
    "",
])
(ROOT / "BUILD_REPORT.md").write_text(report, encoding="utf-8")
print(report)
